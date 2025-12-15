"""
PowerPoint (PPTX) Document Extractor

Extracts text, tables, charts, and metadata from PowerPoint presentations.
"""

from pathlib import Path
from typing import Dict, List, Any, Optional
from io import BytesIO
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract

from ..utils import (
    detect_performance_blocks,
    detect_source_notes,
    detect_language_safe,
    detect_disclaimer_flags,
    detect_glossary,
    detect_legal_notice,
    categorize_disclaimer,
    extract_title_fields,
    extract_identifiers,
    extract_issuer_mentions,
    analyze_performance_sentence,
    parse_table_source,
    detect_countries,
    country_entries,
    normalize_table_rows,
    extract_pptx_table,
    prepare_image_for_ocr,
    extract_text_from_image,
)
from ...parsers.registration import COUNTRY_PATTERNS


def extract_pptx(
    file_path: Path,
    chart_analyzer: Optional[Any] = None,
    enable_chart_analysis: bool = True,
    tesseract_lang: str = "eng",
    tesseract_config: str = "--psm 6"
) -> Dict[str, Any]:
    """
    Extract text and tables from PowerPoint file.
    
    Args:
        file_path: Path to PPTX file
        chart_analyzer: Optional chart analyzer instance
        enable_chart_analysis: Whether to enable chart analysis
        tesseract_lang: Tesseract language code
        tesseract_config: Tesseract configuration
        
    Returns:
        Dict with extracted data
    """
    if Presentation is None:
        raise ImportError("python-pptx not installed")
    
    prs = Presentation(str(file_path))
    
    full_text = []
    slides_data = []
    tables_data = []
    slide_summaries = []
    performance_sections = []
    disclaimer_slides = []
    glossary_slides = []
    country_mentions = set()
    table_sources = []
    identifiers: List[Dict[str, Any]] = []
    performance_table_entries: List[Dict[str, Any]] = []
    country_entries_list: List[Dict[str, Any]] = []
    issuer_mentions: List[Dict[str, Any]] = []
    legal_notice_slide = None
    charts_data: List[Dict[str, Any]] = []
    
    debug_log = open('c:/temp/extraction_debug.log', 'a', encoding='utf-8')
    debug_log.write("\n=== Starting PowerPoint extraction ===\n")
    debug_log.flush()
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_ocr = []
        raw_source_notes = []
        
        debug_log.write(f"Processing slide {slide_idx}, shapes: {len(slide.shapes)}\n")
        debug_log.flush()
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                slide_text.append(text)
                full_text.append(text)
                raw_source_notes.extend(detect_source_notes(text))
            
            # Extract tables
            table_data = None
            try:
                if (hasattr(shape, "shape_type") and 
                    shape.shape_type == MSO_SHAPE_TYPE.TABLE and 
                    hasattr(shape, "table")):
                    table_data = extract_pptx_table(shape.table)
            except (ImportError, AttributeError, Exception):
                pass
            
            if table_data:
                table_index = len(tables_data) + 1
                tables_data.append({
                    'slide_number': slide_idx,
                    'table_index': table_index,
                    'data': table_data
                })
                for entry in normalize_table_rows(table_data):
                    entry.update({'slide_number': slide_idx, 'table_index': table_index})
                    performance_table_entries.append(entry)
 
            # OCR and chart analysis on embedded images
            if (MSO_SHAPE_TYPE is not None
                and Image is not None
                and hasattr(shape, "shape_type")
                and shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
                debug_log.write(f"  Found PICTURE shape on slide {slide_idx}, chart_analyzer enabled: {enable_chart_analysis}\n")
                debug_log.flush()
                try:
                    image_bytes = shape.image.blob
                    
                    # Try chart analysis first (if enabled)
                    if enable_chart_analysis and chart_analyzer:
                        debug_log.write("  Calling chart analyzer...\n")
                        debug_log.flush()
                        try:
                            chart_data = chart_analyzer.extract_chart_data_from_image(
                                image_bytes,
                                location={'slide_number': slide_idx}
                            )
                            debug_log.write(f"  Chart analysis result: is_chart={chart_data.get('is_chart')}, type={chart_data.get('chart_type')}\n")
                            debug_log.flush()
                            
                            # Check if it's a chart (either explicitly marked or has chart data)
                            is_chart = chart_data.get('is_chart', False)
                            has_chart_data = (
                                chart_data.get('chart_type') or
                                len(chart_data.get('data_points', [])) > 0 or
                                len(chart_data.get('performance_values', [])) > 0 or
                                chart_data.get('confidence', 0.0) > 0.5
                            )
                            
                            # If it's a chart or has chart data, process it
                            if is_chart or has_chart_data:
                                # Auto-mark as chart if we have chart data
                                if not is_chart and has_chart_data:
                                    is_chart = True
                                    chart_data['is_chart'] = True
                                
                                # Add chart performance values to performance sections
                                perf_values = chart_data.get('performance_values', [])
                                if perf_values:
                                    if not any(s.get('slide_number') == slide_idx for s in performance_sections):
                                        performance_sections.append({
                                            'slide_number': slide_idx,
                                            'sentences': [],
                                            'entries': []
                                        })
                                    section = next(s for s in performance_sections if s.get('slide_number') == slide_idx)
                                    section['entries'].extend(perf_values)
                                
                                # Add chart source/date info
                                source_info = chart_data.get('source_date_info', {})
                                if source_info.get('has_source') or source_info.get('has_date'):
                                    raw_source_notes.append(
                                        f"Source: {source_info.get('source_text', '')} "
                                        f"Date: {source_info.get('date_text', '')}"
                                    )
                                
                                # Store chart data
                                charts_data.append({
                                    'slide_number': slide_idx,
                                    'is_chart': is_chart,  # Always include is_chart field
                                    'chart_type': chart_data.get('chart_type'),
                                    'chart_title': chart_data.get('chart_title'),
                                    'data_points': chart_data.get('data_points', []),
                                    'performance_values': chart_data.get('performance_values', []),
                                    'source_date_info': source_info,
                                    'confidence': chart_data.get('confidence', 0.0)
                                })
                                debug_log.write(f"  Added chart to charts_data, total charts: {len(charts_data)}\n")
                                debug_log.flush()
                        except Exception as e:
                            debug_log.write(f"  Chart analysis exception: {e}\n")
                            debug_log.write(f"  Exception type: {type(e).__name__}\n")
                            import traceback
                            debug_log.write(f"  Traceback: {traceback.format_exc()}\n")
                            debug_log.flush()
                            print(f"WARNING: Chart analysis failed for slide {slide_idx}: {e}")
                            pass
                    
                    # Also do OCR for text extraction
                    if pytesseract is not None:
                        with Image.open(BytesIO(image_bytes)) as img:
                            img = prepare_image_for_ocr(img)
                            ocr_text = extract_text_from_image(img, tesseract_lang, tesseract_config)
                            if ocr_text:
                                slide_ocr.append(ocr_text)
                                full_text.append(ocr_text)
                                raw_source_notes.extend(detect_source_notes(ocr_text))
                except Exception:
                    continue
        
        combined_text = "\n".join([*slide_text, *slide_ocr])
        # Compile regex patterns for detect_countries
        compiled_patterns = {country: re.compile(pattern, re.IGNORECASE) for country, pattern in COUNTRY_PATTERNS.items()}
        slide_countries = detect_countries(combined_text, compiled_patterns)
        country_mentions.update(slide_countries)

        heading = slide_text[0] if slide_text else None
        country_entries_list.extend(
            country_entries(
                slide_countries,
                heading,
                combined_text,
                {'slide_number': slide_idx}
            )
        )

        issuer_mentions.extend(
            extract_issuer_mentions(
                combined_text,
                {'slide_number': slide_idx, 'heading': heading}
            )
        )

        identifier_entries = extract_identifiers(
            combined_text,
            {'slide_number': slide_idx}
        )
        if identifier_entries:
            identifiers.extend(identifier_entries)

        performance_sentences = detect_performance_blocks(combined_text)
        if performance_sentences:
            entries = [analyze_performance_sentence(sentence) for sentence in performance_sentences]
            performance_sections.append({
                'slide_number': slide_idx,
                'sentences': performance_sentences,
                'entries': entries
            })

        if detect_disclaimer_flags(combined_text):
            disclaimer_slides.append(slide_idx)

        if detect_glossary(combined_text):
            glossary_slides.append(slide_idx)

        if legal_notice_slide is None and detect_legal_notice(combined_text):
            legal_notice_slide = slide_idx

        table_sources.extend({
            'slide_number': slide_idx,
            **parse_table_source(note)
        } for note in raw_source_notes)

        slides_data.append({
            'slide_number': slide_idx,
            'text': '\n'.join(slide_text),
            'ocr_text': '\n'.join(slide_ocr) if slide_ocr else None
        })

        slide_summaries.append({
            'slide_number': slide_idx,
            'is_title_candidate': slide_idx == 1,
            'contains_disclaimer': slide_idx in disclaimer_slides,
            'disclaimer_categories': categorize_disclaimer(combined_text) if slide_idx in disclaimer_slides else [],
            'contains_performance': bool(performance_sentences),
            'contains_glossary': slide_idx in glossary_slides,
            'countries': slide_countries,
            'language': detect_language_safe(combined_text),
        })

        if slide_idx == 1:
            title_info = extract_title_fields(combined_text)

    structure = {
        'title_slide': 1 if slides_data else None,
        'disclaimer_slides': disclaimer_slides,
        'performance_slides': [section['slide_number'] for section in performance_sections],
        'glossary_slides': glossary_slides,
        'countries_detected': sorted(country_mentions),
        'disclaimer_categories': {
            slide['slide_number']: slide['disclaimer_categories']
            for slide in slide_summaries if slide['disclaimer_categories']
        },
        'country_entries': country_entries_list,
        'issuer_mentions': issuer_mentions,
        'has_glossary': bool(glossary_slides),
        'has_management_notice': legal_notice_slide is not None,
        'legal_notice_slide': legal_notice_slide,
    }

    debug_log.write(f"\n=== PowerPoint extraction complete: {len(charts_data)} charts found ===\n")
    debug_log.close()

    return {
        'text': '\n'.join(full_text),
        'slides': slides_data,
        'slide_summaries': slide_summaries,
        'tables': tables_data,
        'total_slides': len(slides_data),
        'total_tables': len(tables_data),
        'table_sources': table_sources,
        'performance_sections': performance_sections,
        'performance_table_entries': performance_table_entries,
        'charts': charts_data,
        'total_charts': len(charts_data),
        'structure': structure,
        'identifiers': identifiers,
        'country_entries': country_entries_list,
        'issuer_mentions': issuer_mentions,
        'title_information': title_info if 'title_info' in locals() else None,
        'ocr': pytesseract is not None and any(slide.get('ocr_text') for slide in slides_data)
    }

