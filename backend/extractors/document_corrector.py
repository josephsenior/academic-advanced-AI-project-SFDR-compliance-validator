"""
Document Corrector

Applies validation fixes to documents based on validation results.
Currently supports PowerPoint (.pptx) format.

Auto-fixes:
- Missing source/date on tables/charts
- Missing disclaimers (when enabled)

Flags for manual review:
- Numerical inconsistencies
- Cross-reference mismatches
"""

from pathlib import Path
from typing import Dict, Any, List, Optional, TYPE_CHECKING
from datetime import datetime
import logging

from .agents.data_consistency_agent import DataConsistencyResult
from .utils.issue_highlighter import IssueHighlighter

logger = logging.getLogger(__name__)

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType
else:
    PresentationType = Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None


class CorrectionResult:
    """Result of document correction operation"""
    
    def __init__(self):
        self.corrected_path: Optional[str] = None
        self.fixes_applied: List[Dict[str, Any]] = []
        self.fixes_failed: List[Dict[str, Any]] = []
        self.manual_review_required: List[Dict[str, Any]] = []
        self.success: bool = False
        self.error_message: Optional[str] = None
        # Visual change tracking
        self.changes_by_slide: Dict[int, List[Dict[str, Any]]] = {}


class DocumentCorrector:
    """
    Apply validation fixes to documents.
    
    Currently supports:
    - PowerPoint (.pptx): Source/date fixes, disclaimer fixes
    """
    
    def __init__(self):
        self.supported_formats = {
            '.pptx': self._correct_pptx,
        }
        self._current_prs = None  # Store current presentation for slide dimension access
        
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        # Attempt to load slide content mapping from project data folder to help map
        # document-wide issues to the most relevant slide.
        try:
            from pathlib import Path
            project_root = Path(__file__).resolve().parents[3]
            mapping_path = project_root / "data" / "slide_content_map.json"
            if mapping_path.exists():
                import json
                self.slide_content_map = {int(k): v for k, v in json.loads(mapping_path.read_text(encoding="utf-8")).items()}
            else:
                self.slide_content_map = {}
        except Exception:
            self.slide_content_map = {}
    
    def correct(
        self,
        original_path: str,
        validation_result: DataConsistencyResult,
        disclaimer_result: Optional[Any] = None,
        output_path: Optional[str] = None,
        **kwargs # Ignore legacy auto_fix arguments
    ) -> CorrectionResult:
        """
        Apply highlighting and suggestions to document.
        
        Args:
            original_path: Path to original document
            validation_result: DataConsistencyResult from validation
            disclaimer_result: Optional disclaimer validation result
            output_path: Path for corrected document
            
        Returns:
            CorrectionResult with details of highlights applied
        """
        result = CorrectionResult()
        input_path_obj = Path(original_path)
        
        if not input_path_obj.exists():
            result.error_message = f"File not found: {input_path_obj}"
            return result
        
        # Determine output path
        if output_path is None:
            stem = input_path_obj.stem
            suffix = input_path_obj.suffix
            output_path_obj = input_path_obj.parent / f"{stem}_corrected{suffix}"
        else:
            output_path_obj = Path(output_path)
        
        # Get file extension
        file_ext = input_path_obj.suffix.lower()
        
        if file_ext not in self.supported_formats:
            result.error_message = f"Unsupported format: {file_ext}. Supported: {list(self.supported_formats.keys())}"
            return result
        
        try:
            # Call appropriate corrector
            corrector_func = self.supported_formats[file_ext]
            # ONLY pass necessary args for highlighting
            corrector_func(input_path_obj, output_path_obj, validation_result, result)
            
            result.corrected_path = str(output_path_obj)
            result.success = True
            
        except Exception as e:
            result.error_message = f"Correction failed: {str(e)}"
            result.success = False
        
        return result
    
    def _correct_pptx(
        self,
        input_path: Path,
        output_path: Path,
        validation_result: DataConsistencyResult,
        result: CorrectionResult
    ) -> None:
        """Apply highlighting to PowerPoint presentation"""
        
        # Load presentation
        prs = Presentation(str(input_path))
        
        # Store prs reference for slide dimension access
        self._current_prs = prs
        
        # Build index maps for deterministic finding of shapes
        self._build_shape_indices(prs)

        # ONLY HIGHLIGHTING - NO AUTO-FIXES
        logger.info("Highlighting document issues (Auto-fixes are disabled)")
        
        # Highlight all issues in the document
        try:
            self._highlight_issues(prs, validation_result, result)
        except Exception as e:
            logger.error(f"Error in _highlight_issues: {e}", exc_info=True)
            result.fixes_failed.append({"issue": "all_highlights", "reason": f"Internal error: {str(e)}"})

        # Flag numerical issues for manual review
        self._flag_numerical_issues(validation_result, result)
        
        # Flag cross-reference issues for manual review
        self._flag_cross_reference_issues(validation_result, result)
        
        # Save presentation
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            logger.info(f"Saving highlighted presentation to {output_path}")
            prs.save(str(output_path))
            logger.info("Successfully saved highlighted presentation")
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}", exc_info=True)
            raise e
        
        # Clean up
        self._current_prs = None

    def _get_issues_from_result(self, result: Any, field_name: str) -> List[Any]:
        """Helper to get issues from validation_result regardless of if it's a dict or model."""
        self._log_debug(f"Retrieving issues for field: {field_name}")
        
        # 1. Handle DICT input
        if isinstance(result, dict):
            inner = result.get("validation_result", result)
            
            # Direct check in inner dict
            if field_name in inner:
                res = inner[field_name]
                if isinstance(res, list):
                    self._log_debug(f"Found {len(res)} issues in dict field '{field_name}'")
                    return res
                
            # Category fallback
            categorized = inner.get("issues_by_category", {})
            mapping = {
                "source_date_issues": "source_date",
                "numerical_inconsistencies": "numerical",
                "cross_reference_issues": "cross_reference"
            }
            
            if field_name in mapping:
                cat_key = mapping[field_name]
                res = categorized.get(cat_key, [])
                self._log_debug(f"Mapped {field_name} to category {cat_key}, found {len(res)} items")
                return res
            
            # Universal fallback for dict: check all common names
            if field_name == "compliance_issues":
                for alt_key in ["compliance_issues", "issues", "violations", "all_issues"]:
                    if alt_key in inner and isinstance(inner[alt_key], list) and inner[alt_key]:
                        self._log_debug(f"Found {len(inner[alt_key])} issues in alternative dict key '{alt_key}'")
                        return inner[alt_key]

        # 2. Handle OBJECT input
        # Try direct attribute access
        res = getattr(result, field_name, None)
        if res is not None and isinstance(res, list):
             self._log_debug(f"Found {len(res)} issues in object attribute '{field_name}'")
             return res
             
        # Fallback for models: check properties or common aliases
        if field_name == "compliance_issues":
             for alt_attr in ["compliance_issues", "issues", "violations"]:
                 res = getattr(result, alt_attr, None)
                 if res is not None and isinstance(res, list) and res:
                     self._log_debug(f"Found {len(res)} issues in alternative attribute '{alt_attr}'")
                     return res

             # Last resort: Combine all legacy fields
             combined = []
             combined.extend(getattr(result, "source_date_issues", []))
             combined.extend(getattr(result, "numerical_inconsistencies", []))
             combined.extend(getattr(result, "cross_reference_issues", []))
             if combined:
                 self._log_debug(f"Combined {len(combined)} issues from legacy attributes")
             return combined

        return []

    def _get_val(self, obj: Any, field: str, default: Any = None) -> Any:
        """Helper to get value from object or dict."""
        if isinstance(obj, dict):
            return obj.get(field, default)
        return getattr(obj, field, default)
    
    def _flag_numerical_issues(
        self,
        validation_result: DataConsistencyResult,
        result: CorrectionResult
    ) -> None:
        """Flag numerical issues for manual review (don't auto-fix)"""
        numerical_inconsistencies = self._get_issues_from_result(validation_result, "numerical_inconsistencies")
        
        for inc in numerical_inconsistencies:
            result.manual_review_required.append({
                "type": "numerical_inconsistency",
                "location": self._get_val(inc, "location"),
                "issue": self._get_val(inc, "message"),
                "document_value": self._get_val(inc, "document_value"),
                "reference_value": self._get_val(inc, "reference_value"),
                "action": "Review and correct manually"
            })
    
    def _flag_cross_reference_issues(
        self,
        validation_result: DataConsistencyResult,
        result: CorrectionResult
    ) -> None:
        """Flag cross-reference issues for manual review"""
        cross_reference_issues = self._get_issues_from_result(validation_result, "cross_reference_issues")
        
        for issue in cross_reference_issues:
            result.manual_review_required.append({
                "type": "cross_reference",
                "location": self._get_val(issue, "location"),
                "issue": self._get_val(issue, "message"),
                "value1": self._get_val(issue, "value1"),
                "value2": self._get_val(issue, "value2"),
                "action": "Review and correct manually"
            })

    def _log_debug(self, msg):
        try:
            # Use absolute path to ensure we can find the log
            log_path = r"C:\Users\GIGABYTE\Desktop\Portfolio\Advanced Ai Project\debug_corrector.txt"
            with open(log_path, "a", encoding='utf-8') as f:
                f.write(f"{datetime.now()} - {msg}\n")
        except Exception:
            pass

    def _highlight_issues(
        self,
        prs: PresentationType,
        validation_result: DataConsistencyResult,
        result: CorrectionResult
    ) -> None:
        """Add visual markers (borders) for issues and add details to Speaker Notes."""
        logger.info(f"!!! STARTING HIGHLIGHTING for {len(prs.slides)} slides")

        # Robustly get issues from result
        compliance_issues = self._get_issues_from_result(validation_result, "compliance_issues")
        
        all_issues = []
        if compliance_issues:
            all_issues = list(compliance_issues)
        else:
             # Combined fallback
             all_issues.extend(self._get_issues_from_result(validation_result, "source_date_issues"))
             all_issues.extend(self._get_issues_from_result(validation_result, "numerical_inconsistencies"))
             all_issues.extend(self._get_issues_from_result(validation_result, "cross_reference_issues"))
        
        logger.info(f"!!! FOUND {len(all_issues)} issues to process")

        issues_by_slide: Dict[int, List[Dict[str, Any]]] = {}
        all_issue_dicts: List[Dict[str, Any]] = []

        for issue in all_issues:
            # Conversion logic...
            if isinstance(issue, dict):
                issue_dict = issue
            elif hasattr(issue, 'model_dump'):
                issue_dict = issue.model_dump()
            elif hasattr(issue, 'dict'):
                issue_dict = issue.dict()
            else:
                issue_dict = {k: getattr(issue, k) for k in dir(issue) if not k.startswith('_') and not callable(getattr(issue, k))}
            
            all_issue_dicts.append(issue_dict)
            
            slide_num = issue_dict.get('slide_number')
            if not slide_num or not isinstance(slide_num, int):
                slide_num = 1 # Fallback
                
            if slide_num >= 1:
                if slide_num not in issues_by_slide:
                    issues_by_slide[slide_num] = []
                issues_by_slide[slide_num].append(issue_dict)

        # Apply slide-specific highlights
        for slide_num, issues in issues_by_slide.items():
            if slide_num > len(prs.slides):
                continue
            
            slide = prs.slides[slide_num - 1]
            logger.info(f"!!! HIGHLIGHTING Slide {slide_num} ({len(issues)} issues)")
            
            # 1. Update Speaker Notes
            IssueHighlighter.add_speaker_notes(slide, issues)

            # 2. Add Visual Highlights
            doc_box_y_offset = 0.5 
            
            for issue in issues:
                target_shape = None
                
                # Check for direct shape identification (preferred)
                table_idx = issue.get('table_index')
                chart_idx = issue.get('chart_index')
                
                if table_idx is not None:
                    target_shape = self._find_shape_by_index(slide, MSO_SHAPE_TYPE.TABLE, table_idx)
                elif chart_idx is not None:
                    target_shape = self._find_shape_by_index(slide, MSO_SHAPE_TYPE.CHART, chart_idx)
                    if not target_shape:
                        GRAPHIC_FRAME = getattr(MSO_SHAPE_TYPE, 'GRAPHIC_FRAME', 14)
                        target_shape = self._find_shape_by_index(slide, GRAPHIC_FRAME, chart_idx)
                
                # Context-based lookup
                if not target_shape:
                    context = issue.get('context')
                    if context:
                        target_shape = self._find_shape_by_text(slide, context)
                
                # Heuristic-based lookup for "exact location" if still not found
                if not target_shape:
                    issue_type = issue.get('issue_type', '').lower()
                    
                    # Performance issues -> find first chart or table
                    if 'performance' in issue_type:
                        for shape in slide.shapes:
                            if (hasattr(shape, "has_chart") and shape.has_chart) or \
                               (hasattr(shape, "has_table") and shape.has_table):
                                target_shape = shape
                                break
                    
                    # Risk/SRI issues -> find shape with "risk" or "sri"
                    elif 'risk' in issue_type or 'sri' in issue_type:
                        target_shape = self._find_shape_by_text(slide, "risk") or \
                                       self._find_shape_by_text(slide, "sri")
                    
                    # Disclaimer issues -> find shape with "disclaimer" or long text
                    elif 'disclaimer' in issue_type:
                        target_shape = self._find_shape_by_text(slide, "disclaimer") or \
                                       self._find_shape_by_text(slide, "investor")
                    
                    # NAV issues
                    elif 'nav' in issue_type:
                        target_shape = self._find_shape_by_text(slide, "nav")
                    
                    # Date/Source issues
                    elif 'source' in issue_type or 'date' in issue_type:
                         target_shape = self._find_shape_by_text(slide, "source") or \
                                        self._find_shape_by_text(slide, "date") or \
                                        self._find_shape_by_text(slide, "202") # year prefix
                    
                    # If still not found, and there is "context", try a fuzzy match
                    if not target_shape and issue.get('context'):
                         target_shape = self._find_shape_by_text(slide, issue.get('context'))

                if target_shape:
                    IssueHighlighter.add_red_border_to_shape(slide, target_shape, issue, border_width=5.0)
                    IssueHighlighter.add_issue_marker(slide, target_shape, issue)
                    IssueHighlighter.add_suggestion_box(slide, target_shape, issue)
                else:
                    # Generic box fallback
                    IssueHighlighter.add_document_issue_box(slide, issue, (6.0, doc_box_y_offset))
                    doc_box_y_offset += 1.4
        
        # No longer adding summary slide as per user request (ODDO BHF requirement)
        pass

    def _find_shape_by_text(self, slide: Any, search_text: str) -> Optional[Any]:
        """Search for a shape on the slide that contains the specified text."""
        if not search_text:
            return None
            
        search_text_lower = search_text.lower().strip()
        
        # Try finding a shape where the text is exactly or clearly part of it
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text:
                continue
                
            shape_text = shape.text.lower()
            if search_text_lower in shape_text:
                return shape
                
        return None

    def _bring_to_front(self, shape: Any) -> None:
        """Ensure a shape is at the top of the Z-order so it is visible."""
        try:
            # In python-pptx, moving to front can be done by re-appending the element
            # or using the shape's Z-order if exposed (it's internal in some versions)
            # Standard way to move to front:
            el = shape.element
            el.getparent().append(el)
        except Exception as e:
            logger.warning(f"Failed to bring shape to front: {e}")

    def _build_shape_indices(self, prs: PresentationType) -> None:
        """Build maps of Global Index -> Shape for deterministic lookup."""
        # Global maps (1-based global index -> shape)
        self.table_map = {}
        self.global_table_idx = 0

        # Per-slide maps: slide_index (1-based) -> { local_table_index (0-based) : shape }
        self.table_map_by_slide = {}
        self.chart_map_by_slide = {}

        for slide_index, slide in enumerate(prs.slides, start=1):
            local_table_idx = 0
            local_chart_idx = 0
            self.table_map_by_slide[slide_index] = {}
            self.chart_map_by_slide[slide_index] = {}

            for shape in slide.shapes:
                # Check for Table
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE and hasattr(shape, "has_table") and shape.has_table:
                    # global index
                    self.global_table_idx += 1
                    self.table_map[self.global_table_idx] = shape
                    # per-slide mapping (use 0-based local index)
                    self.table_map_by_slide[slide_index][local_table_idx] = shape
                    local_table_idx += 1
                # Check for Chart/Graphic frame (count as charts)
                elif (hasattr(shape, "has_chart") and shape.has_chart) or shape.shape_type == MSO_SHAPE_TYPE.CHART or getattr(MSO_SHAPE_TYPE, 'GRAPHIC_FRAME', None) == shape.shape_type:
                    # per-slide chart map
                    self.chart_map_by_slide[slide_index][local_chart_idx] = shape
                    local_chart_idx += 1

    def _find_shape_by_index(self, slide: Any, shape_type: Any, index: int) -> Optional[Any]:
        """Find a shape of a specific type by its index."""
        # Determine slide index if we have a current presentation
        slide_index = None
        try:
            if hasattr(self, "_current_prs") and self._current_prs:
                for i, s in enumerate(self._current_prs.slides, start=1):
                    if s is slide:
                        slide_index = i
                        break
        except Exception:
            slide_index = None

        # For tables, prefer per-slide map if available
        if shape_type == MSO_SHAPE_TYPE.TABLE:
            # Try per-slide map lookup (index is local per-slide index or global index)
            if slide_index and hasattr(self, "table_map_by_slide"):
                # If caller passed a global index, try to map it to local by searching global map first
                # But most callers pass a local index; try local directly
                local_map = self.table_map_by_slide.get(slide_index, {})
                if index in local_map:
                    return local_map[index]
                # Fallback: if index is global, check global map and verify it belongs to this slide
                if hasattr(self, "table_map") and index in self.table_map:
                    candidate = self.table_map[index]
                    try:
                        for s in slide.shapes:
                            if s is candidate:
                                return candidate
                    except Exception:
                        return candidate
            else:
                # No per-slide maps available, fallback to global map behavior
                if hasattr(self, 'table_map') and index in self.table_map:
                    shape = self.table_map[index]
                    try:
                        for s in slide.shapes:
                            if s is shape:
                                return shape
                    except Exception:
                        return shape

        # Fallback for charts or if map failed. Prefer per-slide chart map if available.
        count = 0
        if shape_type in [MSO_SHAPE_TYPE.CHART, getattr(MSO_SHAPE_TYPE, 'GRAPHIC_FRAME', 14)]:
            # Try per-slide chart map if we found a slide_index
            if slide_index and hasattr(self, "chart_map_by_slide"):
                local_chart_map = self.chart_map_by_slide.get(slide_index, {})
                if index in local_chart_map:
                    return local_chart_map[index]

        for shape in slide.shapes:
            # For charts, check has_chart or shape_type
            if shape_type in [MSO_SHAPE_TYPE.CHART, getattr(MSO_SHAPE_TYPE, 'GRAPHIC_FRAME', 14)]:
                if (hasattr(shape, "has_chart") and shape.has_chart) or shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    # Note: This finds N-th chart on slide, but issue might pass Global Index.
                    # Without global chart map (which requires LLM logic), this is best effort.
                    if count == index:
                        return shape
                    count += 1
            elif shape.shape_type == shape_type:
                if count == index:
                    return shape
                count += 1
        return None
