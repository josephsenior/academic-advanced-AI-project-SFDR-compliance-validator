"""
Issue Highlighter Utility

Provides enhanced visual highlighting and annotation capabilities for marking
compliance issues in PowerPoint presentations with red bordered boxes.

Features:
- Color-coded issue markers based on severity
- Contextual suggestion boxes (green)
- Issue legend/summary slide generation
- Speaker notes integration
- Z-order management for visibility
"""

from typing import Dict, Any, List, Tuple
import logging

logger = logging.getLogger(__name__)

try:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_FILL
    from pptx.enum.text import MSO_AUTO_SIZE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class IssueHighlighter:
    """Utility for highlighting compliance issues in presentations."""
    
    # Severity color mapping - ENSURE BRIGHT COLORS
    SEVERITY_COLORS = {
        'critical': RGBColor(255, 0, 0),        # Bright Red
        'error': RGBColor(255, 0, 0),           # Bright Red
        'high': RGBColor(255, 128, 0),         # Bright Orange
        'warning': RGBColor(255, 128, 0),       # Bright Orange
        'info': RGBColor(0, 128, 255),          # Bright Blue
        'low': RGBColor(128, 128, 128),         # Gray
    }
    
    # Background colors for highlighting
    SEVERITY_BG_COLORS = {
        'critical': RGBColor(255, 240, 240),    # Very light red
        'error': RGBColor(255, 245, 245),       # Light red
        'high': RGBColor(255, 245, 230),        # Light orange-red
        'warning': RGBColor(255, 250, 230),     # Light orange
        'info': RGBColor(230, 240, 255),        # Light blue
        'low': RGBColor(240, 240, 240),         # Light gray
    }
    
    @staticmethod
    def get_severity_color(severity: str) -> RGBColor:
        """Get RGB color for severity level."""
        s = severity.lower()
        if s in ['critical', 'error']: return IssueHighlighter.SEVERITY_COLORS['critical']
        if s in ['high', 'warning']: return IssueHighlighter.SEVERITY_COLORS['high']
        return IssueHighlighter.SEVERITY_COLORS.get(s, RGBColor(100, 100, 100))
    
    @staticmethod
    def get_severity_bg_color(severity: str) -> RGBColor:
        """Get background RGB color for severity level."""
        s = severity.lower()
        if s in ['critical', 'error']: return IssueHighlighter.SEVERITY_BG_COLORS['critical']
        if s in ['high', 'warning']: return IssueHighlighter.SEVERITY_BG_COLORS['high']
        return IssueHighlighter.SEVERITY_BG_COLORS.get(s, RGBColor(240, 240, 240))
    
    @staticmethod
    def _get_val(obj: Any, field: str, default: Any = None) -> Any:
        """Helper to get value from object or dict."""
        if isinstance(obj, dict):
            return obj.get(field, default)
        return getattr(obj, field, default)

    @staticmethod
    def add_red_border_to_shape(slide: Any, shape: Any, issue: Any, border_width: float = 4.0) -> None:
        """
        Add a red or color-coded border around a shape.
        
        Args:
            slide: PowerPoint slide object
            shape: Shape to border
            issue: Issue dictionary or object with severity and other metadata
            border_width: Width of border in points
        """
        if not PPTX_AVAILABLE:
            return
        
        try:
            severity = IssueHighlighter._get_val(issue, 'severity', 'warning').lower()
            # FORCE RED BORDER for issue locations as requested by user
            border_color = RGBColor(255, 0, 0) 
            
            # Calculate border dimensions - slightly larger than shape
            # Pt(914400) = 1 inch
            border = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                shape.left - Pt(4),
                shape.top - Pt(4),
                shape.width + Pt(8),
                shape.height + Pt(8)
            )
            
            # CRITICAL FIX: Use background() to ensure it is transparent to content
            # background() matches slide background, which hides content if it's over a chart/table
            # However, no_fill() is often what's needed for a transparent overlay
            try:
                border.fill.background()
            except:
                # Fallback if background() fails
                pass
            border.line.color.rgb = border_color
            border.line.width = Pt(border_width)
            
            # BRING TO FRONT - Mandatory so user sees it over content
            try:
                el = border.element
                el.getparent().append(el)
            except Exception:
                pass
            
        except Exception as e:
            logger.warning(f"Failed to add red border: {e}")

    @staticmethod
    def add_issue_marker(slide: Any, shape: Any, issue: Any) -> None:
        """
        Add severity marker/label to top-left of shape.
        
        Args:
            slide: PowerPoint slide object
            shape: Shape to mark
            issue: Issue dictionary or object
        """
        if not PPTX_AVAILABLE:
            return
        
        try:
            severity = IssueHighlighter._get_val(issue, 'severity', 'warning').lower()
            color = IssueHighlighter.get_severity_color(severity)
            severity_label = severity.upper()
            
            marker_box = slide.shapes.add_textbox(
                shape.left,
                shape.top - Inches(0.35),
                Inches(1.5),
                Inches(0.3)
            )
            
            marker_box.fill.solid()
            marker_box.fill.fore_color.rgb = color
            marker_box.line.color.rgb = color
            marker_box.line.width = Pt(1)
            
            tf = marker_box.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = f"⚠️ {severity_label}"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)  # White text
            
            # Bring to front
            try:
                el = marker_box.element
                el.getparent().append(el)
            except Exception:
                pass
                
        except Exception as e:
            logger.warning(f"Failed to add issue marker: {e}")

    @staticmethod
    def _find_safe_position(slide: Any, left: int, top: int, width: int, height: int) -> Tuple[int, int]:
        """Find a position that doesn't overlap with existing suggestion or issue boxes."""
        current_left = left
        current_top = top
        padding = Inches(0.1)
        
        # Maximum attempts to find a spot
        for _ in range(10):
            overlap = False
            for shape in slide.shapes:
                # Check if it's one of our boxes by checking fill color
                is_our_box = False
                try:
                    if hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                        rgb = shape.fill.fore_color.rgb
                        # Match light green (suggestions) or severity bg colors
                        if rgb == RGBColor(220, 255, 220) or \
                           rgb == RGBColor(255, 230, 230) or \
                           rgb == RGBColor(255, 245, 200) or \
                           rgb == RGBColor(240, 240, 240):
                            is_our_box = True
                except:
                    pass
                
                if is_our_box:
                    # Check for intersection
                    s_left, s_top, s_width, s_height = shape.left, shape.top, shape.width, shape.height
                    
                    if not (current_left + width < s_left or 
                            current_left > s_left + s_width or 
                            current_top + height < s_top or 
                            current_top > s_top + s_height):
                        overlap = True
                        # Move down if there's an overlap
                        current_top = s_top + s_height + padding
                        break
            
            if not overlap:
                break
                
        return current_left, current_top

    @staticmethod
    def add_suggestion_box(slide: Any, shape: Any, issue: Any) -> None:
        """
        Add green suggestion box with action items.
        
        Args:
            slide: PowerPoint slide object
            shape: Shape to annotate
            issue: Issue dictionary or object with suggestion
        """
        if not PPTX_AVAILABLE:
            return
        
        try:
            suggestion = IssueHighlighter._get_val(issue, 'suggestion', 'Review required')
            
            # Position at top-right
            box_width = Inches(3.5)
            box_left = shape.left + shape.width - box_width - Inches(0.2)
            box_top = shape.top - Inches(0.35)
            
            # Fallback position if too high
            if box_top < Inches(0.1):
                box_top = shape.top + shape.height + Inches(0.1)
            
            # COLLISION AVOIDANCE: Ensure boxes don't render above each other
            box_left, box_top = IssueHighlighter._find_safe_position(slide, box_left, box_top, box_width, Inches(0.65))
            
            callout = slide.shapes.add_textbox(box_left, box_top, box_width, Inches(0.65))
            callout.fill.solid()
            callout.fill.fore_color.rgb = RGBColor(220, 255, 220)  # Light green
            callout.line.color.rgb = RGBColor(0, 180, 0)  # Dark green
            callout.line.width = Pt(1.5)
            
            tf = callout.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            # Ensure text is not hidden by setting margins
            tf.margin_top = Pt(5)
            tf.margin_bottom = Pt(5)
            tf.margin_left = Pt(5)
            tf.margin_right = Pt(5)
            
            p = tf.paragraphs[0]
            
            # Remove truncation to ensure full suggestion is visible
            p.text = f"💡 {suggestion}"
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 128, 0) # Brighter green for suggestion text
            
            # Bring to front
            try:
                el = callout.element
                el.getparent().append(el)
            except Exception:
                pass
                
        except Exception as e:
            logger.warning(f"Failed to add suggestion box: {e}")

    @staticmethod
    def add_document_issue_box(slide: Any, issue: Any, position: Tuple[float, float]) -> None:
        """
        Add a standalone red-bordered issue box for document-wide or missing issues.
        
        Args:
            slide: PowerPoint slide object
            issue: Issue dictionary or object
            position: Tuple of (left, top) in inches
        """
        if not PPTX_AVAILABLE:
            return
        
        try:
            severity = IssueHighlighter._get_val(issue, 'severity', 'warning').lower()
            color = IssueHighlighter.get_severity_color(severity)
            bg_color = IssueHighlighter.get_severity_bg_color(severity)
            
            left_in, top_in = position
            
            # Position relative to right edge if it's document-wide
            if left_in > 5.0: # Heuristic for "on the right"
                 try:
                      # Try to get slide width from parent presentation
                      # Robust navigation: Slide -> Part -> Package -> Presentation
                      prs = slide.part.package.presentation
                      if hasattr(prs, 'slide_width'):
                           sw_in = float(prs.slide_width) / 914400.0 # EMUs to Inches
                           left_in = sw_in - 4.3 # 4 inches wide + 0.3 margin
                 except Exception:
                      pass

            # COLLISION AVOIDANCE: Ensure boxes don't render above each other
            box_left, box_top = IssueHighlighter._find_safe_position(
                slide, 
                Inches(left_in), 
                Inches(top_in), 
                Inches(4), 
                Inches(1.2)
            )
            
            # Create main box with RED BORDER
            main_box = slide.shapes.add_textbox(box_left, box_top, Inches(4), Inches(1.2))
            main_box.fill.solid()
            main_box.fill.fore_color.rgb = bg_color
            main_box.line.color.rgb = color
            main_box.line.width = Pt(3)  # Thick border
            
            tf = main_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            tf.margin_bottom = Inches(0.05)
            tf.margin_top = Inches(0.05)
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            
            # Header: Severity and issue type
            issue_type = IssueHighlighter._get_val(issue, 'issue_type', 'UNKNOWN')
            p1 = tf.paragraphs[0]
            p1.text = f"🔴 [{severity.upper()}] {issue_type}"
            p1.font.size = Pt(10)
            p1.font.bold = True
            p1.font.color.rgb = color
            
            # Main content
            suggestion = IssueHighlighter._get_val(issue, 'suggestion')
            message = IssueHighlighter._get_val(issue, 'message', '')
            
            if suggestion:
                p2 = tf.add_paragraph()
                p2.text = f"💡 {suggestion}"
                p2.font.size = Pt(9)
                p2.font.bold = True
                p2.font.color.rgb = RGBColor(0, 128, 0)  # Green color for suggestion text
            
            if message:
                p3 = tf.add_paragraph()
                message_text = str(message)[:40] + "..." if len(str(message)) > 40 else str(message)
                p3.text = f"({message_text})"
                p3.font.size = Pt(8)
                p3.font.italic = True
                p3.font.color.rgb = RGBColor(120, 120, 120)
            
            # Bring to front
            try:
                el = main_box.element
                el.getparent().append(el)
            except Exception:
                pass
                
        except Exception as e:
            logger.warning(f"Failed to add document issue box: {e}")

    @staticmethod
    def create_issues_summary_slide(prs: Any, issues: List[Any]) -> None:
        """
        Create a comprehensive summary slide listing all issues found.
        
        Args:
            prs: Presentation object
            issues: List of all issues found
        """
        if not PPTX_AVAILABLE:
            return
        
        try:
            # Check if we already added a summary (to avoid duplicates on retry)
            if len(prs.slides) > 0:
                 first_slide = prs.slides[0]
                 if hasattr(first_slide, 'shapes'):
                      for s in first_slide.shapes:
                           if hasattr(s, 'text') and "COMPLIANCE ISSUES SUMMARY" in s.text:
                                return # Already exists
            
            # Create blank slide at beginning
            blank_layout = prs.slide_layouts[6]
            summary_slide = prs.slides.add_slide(blank_layout)
            
            # Move to first position
            xml_slides = prs.slides._sldIdLst
            slides_count = len(xml_slides)
            xml_slides.insert(0, xml_slides[slides_count - 1])
            
            # Title
            title_box = summary_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = "🔍 COMPLIANCE ISSUES SUMMARY"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(200, 0, 0)
            
            # Group by severity
            issues_by_severity = {}
            for issue in issues:
                severity = IssueHighlighter._get_val(issue, 'severity', 'info').lower()
                if severity not in issues_by_severity:
                    issues_by_severity[severity] = []
                issues_by_severity[severity].append(issue)
            
            # Add severity summary boxes
            severity_order = ['critical', 'error', 'high', 'warning', 'info', 'low']
            current_top = 1.2
            
            for severity in severity_order:
                if severity not in issues_by_severity:
                    continue
                
                count = len(issues_by_severity[severity])
                color = IssueHighlighter.get_severity_color(severity)
                
                # Summary box
                box = summary_slide.shapes.add_textbox(
                    Inches(0.5), Inches(current_top), Inches(9), Inches(0.6)
                )
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                box.line.color.rgb = color
                box.line.width = Pt(2)
                
                tf = box.text_frame
                tf.margin_left = Inches(0.15)
                p = tf.paragraphs[0]
                p.text = f"● {severity.upper()}: {count} issue{'s' if count != 1 else ''}"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = color
                
                current_top += 0.75
            
            # Add usage instructions
            instructions_box = summary_slide.shapes.add_textbox(
                Inches(0.5), Inches(current_top + 0.3), Inches(9), Inches(2.5)
            )
            tf = instructions_box.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = "HOW TO USE THIS CORRECTED DOCUMENT:"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            
            instructions = [
                "✓ RED BORDERED BOXES mark all compliance issues requiring attention",
                "✓ GREEN suggestion boxes provide specific fixes for each issue",
                "✓ Issue severity is color-coded: Red=Critical/Error, Orange=High/Warning, Blue=Info",
                "✓ Speaker notes on each slide contain detailed context and regulatory requirements",
                "✓ Review each highlighted item and apply the suggested corrections",
            ]
            
            for instruction in instructions:
                p2 = tf.add_paragraph()
                p2.text = instruction
                p2.font.size = Pt(11)
                p2.font.color.rgb = RGBColor(50, 50, 50)
                
        except Exception as e:
            logger.warning(f"Failed to create issues summary slide: {e}")

    @staticmethod
    def add_speaker_notes(slide: Any, issues: List[Any]) -> None:
        """
        Add comprehensive issue details to speaker notes.
        
        Args:
            slide: PowerPoint slide object
            issues: List of issues on this slide
        """
        try:
            if not slide.has_notes_slide:
                slide.notes_slide
            
            notes_frame = slide.notes_slide.notes_text_frame
            
            # Add separator
            if len(notes_frame.text) > 0:
                p = notes_frame.add_paragraph()
                p.text = "=" * 60
            
            # Add issues section
            p = notes_frame.add_paragraph()
            p.text = f"COMPLIANCE ISSUES ({len(issues)} found):"
            p.font.bold = True
            
            for issue in issues:
                severity = IssueHighlighter._get_val(issue, 'severity', 'LOW').upper()
                message = IssueHighlighter._get_val(issue, 'message', IssueHighlighter._get_val(issue, 'issue', 'Unknown'))
                context = IssueHighlighter._get_val(issue, 'context', '')
                suggestion = IssueHighlighter._get_val(issue, 'suggestion', '')
                location = IssueHighlighter._get_val(issue, 'location', '')
                
                p = notes_frame.add_paragraph()
                p.text = f"\n[{severity}] {message}"
                p.level = 0
                
                if context:
                    p2 = notes_frame.add_paragraph()
                    p2.text = f"Context: {context}"
                    p2.level = 1
                
                if suggestion:
                    p3 = notes_frame.add_paragraph()
                    p3.text = f"Action: {suggestion}"
                    p3.level = 1
                
                if location:
                    p4 = notes_frame.add_paragraph()
                    p4.text = f"Location: {location}"
                    p4.level = 1
                    
        except Exception as e:
            logger.warning(f"Failed to add speaker notes: {e}")
