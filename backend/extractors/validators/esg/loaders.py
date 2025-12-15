"""
ESG Compliance Agent Document Loaders

Classes for loading documents and metadata.
"""

import logging
from pathlib import Path
from typing import Dict
import fitz
from pptx import Presentation
from docx import Document
import json

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Loads documents from various formats (PDF, PPTX, DOCX)."""
    
    def __init__(self, file_path: str):
        self.file_path = Path(file_path)
        self.doc_type = self.file_path.suffix.lower()
    
    def load(self) -> Dict:
        """Load document content based on file type."""
        if self.doc_type == '.pdf':
            logger.info(" Chargement PDF...")
            doc = fitz.open(str(self.file_path))
            full_text = [page.get_text("text") for page in doc]
            logger.info(f"[OK] PDF chargé ({len(doc)} pages)")
            return {"full_text": "\n".join(full_text), "images": [], "slides_data": []}

        elif self.doc_type == '.pptx':
            logger.info("[TARGET] Chargement PPTX...")
            prs = Presentation(str(self.file_path))
            full_text = []
            slide_images = []
            slide_data = []

            for i, slide in enumerate(prs.slides):
                slide_title = ""
                slide_text = []
                image_count = 0

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder and shape.placeholder_format.type == 1:
                            slide_title = shape.text.strip()
                        slide_text.append(shape.text)

                    if hasattr(shape, "shape_type") and shape.shape_type == 13 and image_count == 0:
                        image_count += 1
                        slide_images.append({
                            "path": None,
                            "slide_number": i + 1,
                            "slide_title": slide_title or f"Slide {i+1}"
                        })

                full_text.extend(slide_text)
                slide_data.append({
                    "slide_number": i + 1,
                    "title": slide_title or f"Slide {i+1}",
                    "text": "\n".join(slide_text)
                })

            logger.info(f"[OK] PPTX chargé ({len(prs.slides)} slides, {len(slide_images)} images trouvées)")
            return {
                "full_text": "\n".join(full_text),
                "images": slide_images,
                "slides_data": slide_data
            }

        elif self.doc_type == '.docx':
            logger.info("[DOC] Chargement DOCX...")
            doc = Document(str(self.file_path))
            full_text = [p.text for p in doc.paragraphs if p.text.strip()]
            logger.info(f"[OK] DOCX chargé ({len(full_text)} paragraphes)")
            return {"full_text": "\n".join(full_text), "images": [], "slides_data": []}

        else:
            raise ValueError(f"Format non supporté: {self.doc_type}")


class MetadataLoader:
    """Loads metadata from JSON files."""
    
    @staticmethod
    def load(metadata_file: str) -> dict:
        """Load metadata from JSON file."""
        with open(metadata_file, 'r', encoding='utf-8') as f:
            return json.load(f)

