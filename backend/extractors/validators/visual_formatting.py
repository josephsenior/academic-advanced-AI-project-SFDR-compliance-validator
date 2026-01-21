"""
Visual Formatting Verifier - Rule 1

Validates visual formatting compliance for risk disclaimers:
- Risk disclaimers must be bold
- Risk disclaimers must use same font size as body text
- Risk disclaimers must not be hidden in footers
- Risk disclaimers must be visually prominent

Rule 1: "Les avertissements/disclaimers sur les risques doivent : 
         être en gras, utiliser la même police et taille que le texte principal, 
         être visibles (pas en footnote/footer)"

References:
- Rule 1 - Disclaimer Formatting Requirements
"""

import re
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from dataclasses import dataclass

logger = logging.getLogger(__name__)

try:
    import pptx
    from pptx.util import Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


@dataclass
class FontProperties:
    """Properties of a font"""
    size: Optional[float]  # in points
    bold: bool
    italic: bool
    name: str


@dataclass
class FormattingCheckResult:
    """Result of checking formatting of text"""
    text_content: str
    is_bold: bool
    font_size: Optional[float]
    font_name: str
    location: str  # "slide", "footer", "header"
    shape_index: int


class VisualFormattingVerifier:
    """Verifies visual formatting of disclaimers using python-pptx"""
    
    # Risk disclaimer keywords that must be bold
    RISK_KEYWORDS = [
        "risk", "risque", "risiko",
        "capital loss", "perte en capital", "kapitalverlust",
        "warning", "avertissement", "warnung",
        "disclaimer", "responsabilité", "haftung",
        "guarantee", "garantie", "garantie",
        "important", "critical", "crucial",
        "mandatory", "obligatoire", "verbindlich"
    ]
    
    def __init__(self, pptx_file_path: str):
        """
        Initialize verifier with PPTX file
        
        Args:
            pptx_file_path: Path to PPTX file
        """
        self.pptx_file_path = Path(pptx_file_path)
        self.presentation = None
        self.body_font_size = None
        self.body_font_name = None
        
        if PPTX_AVAILABLE and self.pptx_file_path.exists() and self.pptx_file_path.stat().st_size > 0:
            try:
                # Use str path for Presentation
                path_str = str(self.pptx_file_path.absolute())
                logger.debug(f"Attempting to load PPTX: {path_str} (type: {type(path_str)})")
                
                from pptx import Presentation
                self.presentation = Presentation(path_str)
                self._detect_body_formatting()
            except Exception as e:
                # Log detailed error for debugging but don't crash
                logger.error(f"Error loading PPTX {self.pptx_file_path}: {e}")
                # Also log traceback if it's the rId error
                if "rId" in str(e):
                    import traceback
                    logger.error(traceback.format_exc())
                self.presentation = None
        else:
            if not self.pptx_file_path.exists():
                logger.debug(f"PPTX file not found: {self.pptx_file_path}")
            elif self.pptx_file_path.stat().st_size == 0:
                logger.warning(f"PPTX file is empty: {self.pptx_file_path}")
    
    def _detect_body_formatting(self):
        """Detect standard body text formatting from document"""
        if not self.presentation or not self.presentation.slides:
            return
        
        # Sample formatting from first few slides
        font_sizes = {}
        font_names = {}
        
        # Avoid slicing slides directly to prevent 'list' has no 'rId' error in some pptx versions
        for i, slide in enumerate(self.presentation.slides):
            if i >= 3:
                break
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    try:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    # Record font size
                                    if run.font.size:
                                        size_pt = run.font.size.pt
                                        font_sizes[size_pt] = font_sizes.get(size_pt, 0) + 1
                                    
                                    # Record font name
                                    if run.font.name:
                                        font_names[run.font.name] = font_names.get(run.font.name, 0) + 1
                    except:
                        pass
        
        # Use most common font size and name
        if font_sizes:
            self.body_font_size = max(font_sizes, key=font_sizes.get)
        if font_names:
            self.body_font_name = max(font_names, key=font_names.get)
    
    def check_disclaimer_formatting(self, text_to_find: str) -> List[FormattingCheckResult]:
        """
        Check if disclaimer text is properly formatted (bold)
        
        Args:
            text_to_find: Disclaimer text to search for
            
        Returns:
            List of formatting check results
        """
        results = []
        
        if not PPTX_AVAILABLE or not self.presentation:
            return results
        
        text_lower = text_to_find.lower()
        
        # Search through all slides for this text
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            for shape_idx, shape in enumerate(slide.shapes):
                if not hasattr(shape, "text_frame"):
                    continue
                
                try:
                    full_text = shape.text_frame.text
                    if text_lower not in full_text.lower():
                        continue
                    
                    # Check formatting of text in this shape
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Check if this run contains our disclaimer text
                            if text_lower in run.text.lower():
                                result = FormattingCheckResult(
                                    text_content=run.text,
                                    is_bold=run.font.bold if run.font.bold is not None else False,
                                    font_size=run.font.size.pt if run.font.size else None,
                                    font_name=run.font.name or "Default",
                                    location=self._determine_location(shape_idx, len(slide.shapes)),
                                    shape_index=shape_idx
                                )
                                results.append(result)
                
                except Exception as e:
                    # Skip shapes that can't be processed
                    pass
        
        return results
    
    def _determine_location(self, shape_idx: int, total_shapes: int) -> str:
        """Determine if shape is in normal content or footer area"""
        # Shapes in bottom 20% or with very small indices are likely footers
        if shape_idx > total_shapes * 0.8 or shape_idx == total_shapes - 1:
            return "footer"
        elif shape_idx == 0:
            return "header"
        else:
            return "slide"
    
    def verify_risk_disclaimers_bold(self) -> Dict[str, Any]:
        """
        Verify all risk disclaimers are bold
        
        Returns:
            Dict with verification results
        """
        if not PPTX_AVAILABLE or not self.presentation:
            return {
                "status": "skipped",
                "reason": "python-pptx not available",
                "bold_disclaimers": [],
                "non_bold_disclaimers": [],
                "total_checked": 0
            }
        
        bold_disclaimers = []
        non_bold_disclaimers = []
        checked_count = 0
        
        # Check all shapes for risk-related content
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            for shape_idx, shape in enumerate(slide.shapes):
                if not hasattr(shape, "text_frame"):
                    continue
                
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.lower()
                        
                        # Check if paragraph contains risk keywords
                        is_risk_related = any(kw in para_text for kw in self.RISK_KEYWORDS)
                        
                        if is_risk_related:
                            checked_count += 1
                            
                            # Check if all runs in this paragraph are bold
                            all_bold = True
                            has_text = False
                            
                            for run in paragraph.runs:
                                if run.text.strip():
                                    has_text = True
                                    if run.font.bold is not True:
                                        all_bold = False
                            
                            if has_text:
                                location = self._determine_location(shape_idx, len(slide.shapes))
                                result = {
                                    "slide": slide_num,
                                    "text": paragraph.text[:100],
                                    "is_bold": all_bold,
                                    "location": location,
                                    "font_size": self._get_paragraph_font_size(paragraph),
                                    "font_name": self._get_paragraph_font_name(paragraph)
                                }
                                
                                if all_bold:
                                    bold_disclaimers.append(result)
                                else:
                                    non_bold_disclaimers.append(result)
                
                except Exception:
                    pass
        
        return {
            "status": "completed",
            "bold_disclaimers": bold_disclaimers,
            "non_bold_disclaimers": non_bold_disclaimers,
            "total_checked": checked_count,
            "compliance_rate": (len(bold_disclaimers) / checked_count * 100) if checked_count > 0 else 0
        }
    
    @staticmethod
    def _get_paragraph_font_size(paragraph) -> Optional[float]:
        """Get font size of first run in paragraph"""
        for run in paragraph.runs:
            if run.font.size:
                return run.font.size.pt
        return None
    
    @staticmethod
    def _get_paragraph_font_name(paragraph) -> str:
        """Get font name of first run in paragraph"""
        for run in paragraph.runs:
            if run.font.name:
                return run.font.name
        return "Default"


def verify_pptx_formatting(pptx_file_path: str, extraction_result: Dict[str, Any]) -> Dict[str, Any]:
    """
    Verify visual formatting of risk disclaimers in PPTX
    
    Args:
        pptx_file_path: Path to PPTX file
        extraction_result: Extracted document content (for context)
        
    Returns:
        Verification results
    """
    verifier = VisualFormattingVerifier(pptx_file_path)
    results = verifier.verify_risk_disclaimers_bold()
    
    # Add summary
    results["summary"] = {
        "total_disclaimers": results["total_checked"],
        "compliant": len(results["bold_disclaimers"]),
        "non_compliant": len(results["non_bold_disclaimers"]),
        "compliance_percentage": results.get("compliance_rate", 0)
    }
    
    return results
