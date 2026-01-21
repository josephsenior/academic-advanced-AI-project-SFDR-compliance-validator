"""
Visual Prominence Validator

Validates that disclaimer text has sufficient contrast ratio for accessibility.
Checks disclaimer text color vs background color to ensure WCAG AA compliance
(minimum 4.5:1 contrast ratio for text).

Uses color analysis from PPTX to detect contrast violations.
"""

import logging
import re
from typing import Dict, Any, List, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path
from backend.extractors.rules.models import ComplianceIssue
from backend.extractors.rules.enums import ComplianceIssueType
from .base import BaseValidator

logger = logging.getLogger(__name__)


@dataclass
class ContrastAnalysis:
    """Result of contrast analysis for text element"""
    text: str
    text_color_rgb: Tuple[int, int, int]
    background_color_rgb: Tuple[int, int, int]
    contrast_ratio: float
    wcag_aa_compliant: bool  # >= 4.5:1
    wcag_aaa_compliant: bool  # >= 7:1
    is_violation: bool


class VisualProminenceValidator(BaseValidator):
    """
    Validates visual prominence of disclaimers using contrast analysis.
    
    Checks that disclaimer text has sufficient contrast vs background
    to meet WCAG AA accessibility standards (4.5:1 minimum).
    """
    
    def __init__(self, wcag_standard: str = 'AA'):
        """
        Initialize validator.
        
        Args:
            wcag_standard: 'AA' (4.5:1) or 'AAA' (7:1) minimum contrast
        """
        self.wcag_standard = wcag_standard
        self.min_contrast = 7.0 if wcag_standard == 'AAA' else 4.5
    
    def validate(
        self,
        extraction_result: Dict[str, Any],
        metadata: Optional[Dict[str, Any]] = None,
        client_type: Optional[Any] = None,
        fund_type: Optional[Any] = None
    ) -> List[ComplianceIssue]:
        """
        Validate visual prominence of disclaimer text.
        
        Returns ComplianceIssue for each disclaimer with insufficient contrast.
        """
        issues: List[ComplianceIssue] = []
        
        # Get PPTX file path
        pptx_path = self._get_pptx_path(extraction_result, metadata)
        
        if not pptx_path or not Path(pptx_path).exists():
            logger.debug("No PPTX file found for visual prominence check")
            return issues
        
        try:
            # Analyze contrast in PPTX
            analyses = self._analyze_pptx_contrast(pptx_path)
            
            # Create issues for violations
            for analysis in analyses:
                if analysis.is_violation:
                    issue = ComplianceIssue(
                        issue_type=ComplianceIssueType.VISUAL_CONTRAST_VIOLATION,
                        issue_category='compliance',
                        severity='warning',
                        rule_reference="WCAG AA Accessibility (Visual Prominence)",
                        location=f"Disclaimer text: {analysis.text[:30]}...",
                        message=(
                            f"Disclaimer text contrast ratio {analysis.contrast_ratio:.2f}:1 "
                            f"below WCAG {self.wcag_standard} minimum of {self.min_contrast}:1"
                        ),
                        context=(
                            f"Text: '{analysis.text[:100]}...'\n"
                            f"Text color: RGB{analysis.text_color_rgb}\n"
                            f"Background color: RGB{analysis.background_color_rgb}\n"
                            f"Current ratio: {analysis.contrast_ratio:.2f}:1"
                        ),
                        suggestion=(
                            f"Increase contrast to at least {self.min_contrast}:1. "
                            f"Darken text or lighten background. "
                            f"For readability, use dark text (RGB 0,0,0) on light background (RGB 255,255,255) = 21:1 contrast"
                        ),
                        auto_fixable=False
                    )
                    issues.append(issue)
        
        except Exception as e:
            logger.warning(f"Failed to analyze PPTX contrast: {e}")
        
        return issues
    
    def _get_pptx_path(self, extraction_result: Dict[str, Any], metadata: Optional[Dict[str, Any]]) -> Optional[str]:
        """Extract PPTX file path from extraction result or metadata."""
        # Check metadata
        if metadata and 'file_path' in metadata:
            path = metadata['file_path']
            if isinstance(path, str) and path.endswith('.pptx'):
                return path
        
        # Check extraction result
        if 'file_path' in extraction_result and extraction_result['file_path'].endswith('.pptx'):
            return extraction_result['file_path']
        
        if 'filename' in extraction_result and extraction_result['filename'].endswith('.pptx'):
            return extraction_result['filename']
        
        return None
    
    def _analyze_pptx_contrast(self, pptx_path: str) -> List[ContrastAnalysis]:
        """
        Analyze contrast in PPTX file.
        
        Returns list of ContrastAnalysis for each disclaimer found.
        """
        try:
            from pptx import Presentation
            from pptx.util import Pt, Inches
        except ImportError:
            logger.warning("python-pptx not available for contrast analysis")
            return []
        
        analyses = []
        
        try:
            prs = Presentation(pptx_path)
            
            # Risk/disclaimer keywords to look for
            risk_keywords = [
                'risk', 'warning', 'disclaimer', 'important', 'capital loss',
                'fluctuation', 'performance', 'past', 'guarantee', 'loss',
                'volatility', 'liquidity', 'market', 'credit'
            ]
            
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if not hasattr(shape, 'text_frame'):
                        continue
                    
                    text = shape.text
                    
                    # Check if this looks like a disclaimer
                    if not any(keyword.lower() in text.lower() for keyword in risk_keywords):
                        continue
                    
                    # Analyze each run's contrast
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if not run.text.strip():
                                continue
                            
                            # Extract colors
                            text_color = self._get_text_color(run)
                            bg_color = self._get_background_color(shape, slide, paragraph)
                            
                            if text_color and bg_color:
                                contrast = self._calculate_contrast_ratio(text_color, bg_color)
                                is_violation = contrast < self.min_contrast
                                
                                analysis = ContrastAnalysis(
                                    text=run.text[:100],
                                    text_color_rgb=text_color,
                                    background_color_rgb=bg_color,
                                    contrast_ratio=contrast,
                                    wcag_aa_compliant=contrast >= 4.5,
                                    wcag_aaa_compliant=contrast >= 7.0,
                                    is_violation=is_violation
                                )
                                
                                if is_violation:
                                    analyses.append(analysis)
        
        except Exception as e:
            logger.warning(f"Error analyzing PPTX contrast: {e}")
        
        return analyses
    
    def _get_text_color(self, run) -> Optional[Tuple[int, int, int]]:
        """
        Extract text color from run.
        
        Returns: RGB tuple or None
        """
        try:
            if not hasattr(run.font, 'color'):
                return None
            
            color = run.font.color
            
            if hasattr(color, 'rgb') and color.rgb is not None:
                # RGB color
                rgb = color.rgb
                return (int(rgb[0]), int(rgb[1]), int(rgb[2]))
            elif hasattr(color, 'theme_color'):
                # Theme color - use default text color (usually black)
                return (0, 0, 0)
            else:
                # Default to black
                return (0, 0, 0)
        
        except Exception as e:
            logger.debug(f"Could not extract text color: {e}")
            return (0, 0, 0)  # Default to black
    
    def _get_background_color(self, shape, slide, paragraph) -> Optional[Tuple[int, int, int]]:
        """
        Extract background color (from shape fill or slide background).
        
        Returns: RGB tuple or None
        """
        try:
            # Try shape fill first
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                    rgb = fill.fore_color.rgb
                    return (int(rgb[0]), int(rgb[1]), int(rgb[2]))
            
            # Try slide background
            if hasattr(slide.background, 'fill'):
                fill = slide.background.fill
                if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                    rgb = fill.fore_color.rgb
                    return (int(rgb[0]), int(rgb[1]), int(rgb[2]))
            
            # Default to white background
            return (255, 255, 255)
        
        except Exception as e:
            logger.debug(f"Could not extract background color: {e}")
            return (255, 255, 255)  # Default to white
    
    def _calculate_contrast_ratio(self, fg_rgb: Tuple[int, int, int], bg_rgb: Tuple[int, int, int]) -> float:
        """
        Calculate WCAG contrast ratio between two colors.
        
        Formula: (L1 + 0.05) / (L2 + 0.05) where L is relative luminance
        
        Returns: Contrast ratio (1.0 to 21.0)
        """
        def get_luminance(rgb: Tuple[int, int, int]) -> float:
            """Calculate relative luminance of color."""
            r, g, b = rgb
            # Normalize to 0-1
            r = r / 255.0
            g = g / 255.0
            b = b / 255.0
            
            # Apply gamma correction
            r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
            g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
            b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
            
            # Calculate luminance
            return 0.2126 * r + 0.7152 * g + 0.0722 * b
        
        l1 = get_luminance(fg_rgb)
        l2 = get_luminance(bg_rgb)
        
        # Calculate contrast ratio (lighter / darker)
        lighter = max(l1, l2)
        darker = min(l1, l2)
        
        if lighter + darker < 0.1:  # Both colors very dark
            return 1.0
        
        contrast = (lighter + 0.05) / (darker + 0.05)
        
        # Clamp to valid range
        return min(21.0, max(1.0, contrast))
