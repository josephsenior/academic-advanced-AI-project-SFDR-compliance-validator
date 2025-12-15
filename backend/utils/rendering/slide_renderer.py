"""
Slide Renderer Utility

Converts PowerPoint slides to images for preview and visualization.
"""

from pathlib import Path
from typing import List, Optional

try:
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    HAS_DEPS = True
except ImportError:
    HAS_DEPS = False
    Presentation = None
    Image = None

# Defer comtypes import to runtime rendering functions to avoid import-not-found
HAS_COMTYPES = False


def render_slides_to_images(
    pptx_path: str,
    output_dir: Optional[str] = None,
    scale: float = 1.0
) -> List[str]:
    """
    Convert PowerPoint slides to images.
    
    Args:
        pptx_path: Path to PowerPoint file
        output_dir: Directory to save images (default: same as PPTX)
        scale: Scale factor for images (default: 1.0)
    
    Returns:
        List of image file paths
    """
    if not HAS_DEPS:
        raise ImportError("Required dependencies not available: python-pptx, Pillow")
    
    pptx_path_obj = Path(pptx_path)
    if not pptx_path_obj.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")
    
    # Ensure output_dir is always a Path object
    if output_dir is None:
        output_dir_path: Path = pptx_path_obj.parent / f"{pptx_path_obj.stem}_slides"
    else:
        output_dir_path = Path(output_dir)
    
    output_dir_path.mkdir(parents=True, exist_ok=True)
    
    prs = Presentation(str(pptx_path_obj))
    image_paths = []
    
    # For now, create placeholder images
    # In production, you'd use comtypes (Windows) or libreoffice (Linux) for actual rendering
    for i, slide in enumerate(prs.slides):
        # Create a simple placeholder image
        # In production, use actual slide rendering
        img_path = output_dir_path / f"slide_{i+1:03d}.png"
        create_slide_placeholder(slide, img_path, i+1, scale)
        image_paths.append(str(img_path))
    
    return image_paths


def create_slide_placeholder(slide, output_path: Path, slide_num: int, scale: float = 1.0):
    """Create a placeholder image for a slide (fallback when rendering unavailable)"""
    if not Image:
        raise ImportError("PIL/Pillow not available")
    
    # Standard slide dimensions: 10x7.5 inches at 96 DPI = 960x720 pixels
    width = int(960 * scale)
    height = int(720 * scale)
    
    img = Image.new('RGB', (width, height), color='white')
    draw = ImageDraw.Draw(img)
    
    # Draw slide border
    draw.rectangle((10, 10, width-10, height-10), outline='#cccccc', width=2)
    
    # Add slide number
    try:
        font = ImageFont.truetype("arial.ttf", 24)
    except:
        font = ImageFont.load_default()
    
    draw.text((20, 20), f"Slide {slide_num}", fill='#666666', font=font)
    
    # Extract and display slide text (first 200 chars)
    slide_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            slide_text.append(shape.text[:100])
    
    text_content = "\n".join(slide_text[:5])[:200]
    if text_content:
        y = 60
        for line in text_content.split('\n')[:10]:
            draw.text((30, y), line, fill='#333333', font=font)
            y += 30
            if y > height - 40:
                break
    
    img.save(output_path, 'PNG')
    return output_path


def annotate_slide_image(
    image_path: str,
    annotations: List[dict],
    output_path: Optional[str] = None
) -> str:
    """
    Add visual annotations to a slide image with detailed change markers.
    
    Args:
        image_path: Path to slide image
        annotations: List of annotation dicts with:
            - type: "added", "modified", "issue"
            - x, y: Position (0-1 normalized)
            - text: Annotation text
            - location: Specific location description
            - bbox: Optional bounding box [x1, y1, x2, y2] in normalized coordinates
        output_path: Output path (default: adds _annotated suffix)
    
    Returns:
        Path to annotated image
    """
    if not Image:
        raise ImportError("PIL/Pillow not available")
    
    img = Image.open(image_path)
    draw = ImageDraw.Draw(img)
    
    # Color mapping
    colors = {
        "added": "#48bb78",      # Green
        "modified": "#ed8936",   # Orange
        "issue": "#f56565",      # Red
        "source_date": "#4299e1" # Blue
    }
    
    # Sort annotations by y position (top to bottom)
    sorted_anns = sorted(annotations, key=lambda a: a.get('y', 0.5))
    
    for idx, ann in enumerate(sorted_anns):
        change_type = ann.get('type', ann.get('change_type', 'issue'))
        color = colors.get(change_type, '#666666')
        
        # Handle bounding box if provided
        if 'bbox' in ann and ann['bbox']:
            bbox = ann['bbox']
            x1 = int(bbox[0] * img.width)
            y1 = int(bbox[1] * img.height)
            x2 = int(bbox[2] * img.width)
            y2 = int(bbox[3] * img.height)
            
            # Draw rectangle outline
            draw.rectangle((x1, y1, x2, y2), outline=color, width=3)
            
            # Draw corner markers
            corner_size = 10
            # Top-left
            draw.rectangle((x1, y1, x1+corner_size, y1+3), fill=color)
            draw.rectangle((x1, y1, x1+3, y1+corner_size), fill=color)
            # Top-right
            draw.rectangle((x2-corner_size, y1, x2, y1+3), fill=color)
            draw.rectangle((x2-3, y1, x2, y1+corner_size), fill=color)
            # Bottom-left
            draw.rectangle((x1, y2-3, x1+corner_size, y2), fill=color)
            draw.rectangle((x1, y2-corner_size, x1+3, y2), fill=color)
            # Bottom-right
            draw.rectangle((x2-corner_size, y2-3, x2, y2), fill=color)
            draw.rectangle((x2-3, y2-corner_size, x2, y2), fill=color)
            
            # Position text near the change
            text_x = x1
            text_y = y1 - 30 if y1 > 30 else y2 + 5
        else:
            # Use point-based annotation
            x = int(ann.get('x', 0.5) * img.width)
            y = int(ann.get('y', 0.5) * img.height)
            
            # Draw marker with number
            radius = 18
            draw.ellipse((x-radius, y-radius, x+radius, y+radius), fill=color, outline='white', width=3)
            
            # Draw number inside circle
            try:
                font_num = ImageFont.truetype("arial.ttf", 14)
            except:
                font_num = ImageFont.load_default()
            
            num_text = str(idx + 1)
            bbox_num = draw.textbbox((0, 0), num_text, font=font_num)
            text_width = bbox_num[2] - bbox_num[0]
            text_height = bbox_num[3] - bbox_num[1]
            draw.text((x - text_width//2, y - text_height//2), num_text, fill='white', font=font_num)
            
            text_x = x + radius + 10
            text_y = y - 15
        
        # Draw detailed annotation text
        text = ann.get('text', ann.get('description', ''))
        location = ann.get('location', '')
        
        if text or location:
            try:
                font = ImageFont.truetype("arial.ttf", 11)
                font_bold = ImageFont.truetype("arialbd.ttf", 11)
            except:
                font = ImageFont.load_default()
                font_bold = font
            
            # Build annotation text
            lines = []
            if text:
                lines.append(text)
            if location:
                lines.append(f"Location: {location}")
            
            # Calculate text box size
            max_width = 0
            total_height = 0
            for line in lines:
                bbox = draw.textbbox((0, 0), line, font=font)
                max_width = max(max_width, bbox[2] - bbox[0])
                total_height += (bbox[3] - bbox[1]) + 4
            
            # Draw text background with border
            padding = 8
            text_bg_x1 = text_x - padding
            text_bg_y1 = text_y - padding
            text_bg_x2 = text_x + max_width + padding
            text_bg_y2 = text_y + total_height + padding
            
            # Ensure text box is within image bounds
            if text_bg_x2 > img.width:
                text_bg_x1 = img.width - max_width - padding * 2
                text_bg_x2 = img.width
            if text_bg_y2 > img.height:
                text_bg_y1 = img.height - total_height - padding * 2
                text_bg_y2 = img.height
            
            # Draw background
            draw.rectangle((text_bg_x1, text_bg_y1, text_bg_x2, text_bg_y2), 
                         fill='white', outline=color, width=2)
            
            # Draw text lines
            current_y = text_y
            for i, line in enumerate(lines):
                draw.text((text_x, current_y), line, fill=color, font=font_bold if i == 0 else font)
                bbox = draw.textbbox((text_x, current_y), line, font=font)
                current_y += (bbox[3] - bbox[1]) + 4
    
    if output_path is None:
        output_path = str(Path(image_path).with_suffix('.annotated.png'))
    
    img.save(output_path)
    return output_path

