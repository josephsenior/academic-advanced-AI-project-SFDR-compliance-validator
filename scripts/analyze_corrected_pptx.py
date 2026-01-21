#!/usr/bin/env python
"""Analyze corrected PPTX to verify suggestion box placement."""

from pptx import Presentation
from pathlib import Path
import json

def analyze_pptx(path: Path):
    """Analyze PPTX and count suggestion boxes."""
    prs = Presentation(str(path))
    
    results = {
        "filename": path.name,
        "total_slides": len(prs.slides),
        "slides": []
    }
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "slide_number": slide_idx,
            "total_shapes": len(slide.shapes),
            "suggestion_boxes": [],
            "has_notes": slide.has_notes_slide
        }
        
        # Find suggestion boxes
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text
                    if "Suggested" in text or "💡" in text or "Issue" in text:
                        # Extract first 150 chars
                        snippet = text.replace("\n", " ")[:150]
                        slide_data["suggestion_boxes"].append({
                            "text": snippet,
                            "position": f"({shape.left}, {shape.top})"
                        })
            except Exception:
                continue
        
        results["slides"].append(slide_data)
    
    # Summary
    total_suggestions = sum(len(s["suggestion_boxes"]) for s in results["slides"])
    print(f"\n{'='*60}")
    print(f"Corrected PPTX Analysis: {path.name}")
    print(f"{'='*60}")
    print(f"Total slides: {results['total_slides']}")
    print(f"Total suggestion boxes: {total_suggestions}")
    print("\nBreakdown by slide:")
    
    for slide_data in results["slides"]:
        count = len(slide_data["suggestion_boxes"])
        print(f"\n  Slide {slide_data['slide_number']}: {count} suggestion box(es)")
        for idx, box in enumerate(slide_data["suggestion_boxes"], start=1):
            print(f"    [{idx}] {box['text']}")
    
    print(f"\n{'='*60}\n")
    
    return results

if __name__ == "__main__":
    pptx_path = Path("corrected_documents/5432ef2e-c9e1-4569-b42a-9153c2692419_corrected.pptx")
    
    if not pptx_path.exists():
        print(f"ERROR: File not found: {pptx_path}")
        print("Searching for corrected files...")
        for p in Path("corrected_documents").glob("*.pptx"):
            print(f"  Found: {p}")
        exit(1)
    
    results = analyze_pptx(pptx_path)
    
    # Save results
    output_file = Path("outputs/pptx_analysis.json")
    output_file.parent.mkdir(exist_ok=True)
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(results, f, indent=2, ensure_ascii=False)
    
    print(f"Detailed results saved to: {output_file}")
