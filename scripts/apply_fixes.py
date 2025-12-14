"""Apply simple automated fixes to a PPTX based on a validation report.
This script appends a slide with standard disclaimers and notes addressing
high-priority issues (missing disclaimers, promotional mention, SRI note, etc.).

Usage: python scripts/apply_fixes.py <input_pptx> <output_pptx>
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt

TEMPLATE_TEXT = [
    "Auto-added fixes and disclaimers:",
    "- Capital is not guaranteed and you may not get back the amount you invested.",
    "- Past performance is not a reliable indicator of future results.",
    "- Simulations do not constitute a forecast of future performance.",
    "- Promotional mention: Document promotionnel / Marketing communication.",
    "- Target audience: Professional Client (confirm and adjust as needed).",
    "- Risk Scale (SRI) missing: add SRI (1-7) on slide 2.",
    "Notes: Please review slide 1 for performance content and move it if needed."
]


def add_disclaimer_slide(prs: Presentation):
    # Choose a slide layout that contains title + content if available
    layout = None
    if len(prs.slide_layouts) > 1:
        layout = prs.slide_layouts[1]
    else:
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)
    try:
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
    except Exception:
        # create a textbox if placeholders are not present
        title = None
        body = None

    if title:
        title.text = "Automated Fixes & Disclaimers (review required)"
    else:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        p = tx.text_frame.add_paragraph()
        p.text = "Automated Fixes & Disclaimers (review required)"
        p.font.size = Pt(28)

    if body:
        tf = body.text_frame
        tf.clear()
        for i, line in enumerate(TEMPLATE_TEXT):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(14)
    else:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
        tf = tb.text_frame
        for line in TEMPLATE_TEXT:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python scripts/apply_fixes.py <input_pptx> <output_pptx>")
        sys.exit(2)

    inp = sys.argv[1]
    out = sys.argv[2]

    prs = Presentation(inp)
    add_disclaimer_slide(prs)
    prs.save(out)
    print(f"Saved corrected file to {out}")
