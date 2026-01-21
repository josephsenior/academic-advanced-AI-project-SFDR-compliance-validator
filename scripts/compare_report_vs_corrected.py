import json
import sys
from pathlib import Path
from pptx import Presentation

def find_corrected_pptx(report_filename: str) -> Path | None:
    # Try common corrected locations
    candidates = list(Path("corrected_documents").glob(f"*{report_filename}*")) + list(Path("corrected_documents").glob(f"*{report_filename.replace(' ', '_')}*"))
    if candidates:
        return candidates[0]
    # fallback: search for files containing key tokens
    token = report_filename.split("_")[0] if "_" in report_filename else report_filename[:10]
    for p in Path("corrected_documents").glob("*.pptx"):
        if token in p.name:
            return p
    return None

def extract_texts_from_slide(slide):
    texts = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        except Exception:
            pass
    # notes
    try:
        if slide.has_notes_slide:
            nt = slide.notes_slide.notes_text_frame.text
            if nt:
                texts.append(nt)
    except Exception:
        pass
    return " \n ".join(texts)

def main(report_path: Path, corrected_override: Path | None = None):
    report = json.loads(report_path.read_text(encoding="utf-8"))
    report_filename = report.get("filename") or report_path.stem
    if corrected_override:
        corrected_pptx = corrected_override
    else:
        corrected_pptx = find_corrected_pptx(report_filename)
    if not corrected_pptx:
        print("No corrected PPTX found for", report_filename)
        sys.exit(2)

    print("Using corrected PPTX:", corrected_pptx)
    prs = Presentation(str(corrected_pptx))

    issues = report.get("compliance_issues") or report.get("issues_by_category", {})
    # Normalize to list of issues
    if isinstance(issues, dict):
        # flatten categories
        flat = []
        for cat, lst in issues.items():
            flat.extend(lst)
        issues = flat

    summary = []
    for idx, issue in enumerate(issues, start=1):
        slide_num = issue.get("slide_number")
        suggestion = issue.get("suggestion") or issue.get("message") or ""
        issue_type = issue.get("issue_type") or issue.get("category")
        found = False
        context = ""
        if slide_num and isinstance(slide_num, int) and 1 <= slide_num <= len(prs.slides):
            text = extract_texts_from_slide(prs.slides[slide_num - 1])
            if suggestion and suggestion[:30] in text:
                found = True
                context = f"Found suggestion text on slide {slide_num}"
            else:
                # look for hint keywords like 'Suggested' or '💡' or part of message
                if "Suggested" in text or "💡" in text or any(part in text for part in [issue.get("message","")[:40], issue.get("location","")]):
                    found = True
                    context = f"Found highlight/note on slide {slide_num}"
        else:
            # document-wide: check last slide and notes
            last = prs.slides[-1]
            text = extract_texts_from_slide(last)
            if suggestion and suggestion[:30] in text:
                found = True
                context = "Found suggestion text on last slide"
            elif "Disclaimer" in text or "disclaimer" in text or "Disclaimer:" in text or "The team is subject to change" in text:
                found = True
                context = "Found disclaimer-like text on last slide"

        summary.append({"index": idx, "issue_type": issue_type, "slide_number": slide_num, "found": found, "context": context})

    # Print concise report
    total = len(summary)
    matched = sum(1 for s in summary if s["found"])
    print(f"Checked {total} issues — matched in corrected PPTX: {matched}/{total}")
    for s in summary:
        status = "OK" if s["found"] else "MISSING"
        print(f"[{status}] Issue #{s['index']} type={s['issue_type']} slide={s['slide_number']} — {s['context']}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: compare_report_vs_corrected.py <validation_report.json> [corrected_pptx_path]")
        sys.exit(1)
    report = Path(sys.argv[1])
    corrected = Path(sys.argv[2]) if len(sys.argv) > 2 else None
    main(report, corrected)

