
import sys
import logging
from pathlib import Path

# Add project root to path
sys.path.append(str(Path.cwd()))

from backend.extractors.document_corrector import DocumentCorrector

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def debug_highlighting():
    corrector = DocumentCorrector()
    
    input_path = "dataset/example_1/FINAL 47861-6PG-FR-ODDO BHF Algo Trend US-20250831.pptx"
    output_path = "dataset/example_1/debug_corrected.pptx"
    
    # Mock data consistency result with explicit slide numbers
    validation_result = {
        "compliance_issues": [
            {
                "issue_type": "missing_sri_risk_scale",
                "message": "DEBUG: Missing SRI Risk Scale",
                "severity": "high",
                "location": "Slide 2",
                "slide_number": 2, 
                "suggestion": "Add SRI scale"
            },
            {
                "issue_type": "performance_starts_document",
                "message": "DEBUG: Performance on Slide 1",
                "severity": "error",
                "location": "Slide 1",
                "slide_number": 1, 
                "suggestion": "Move to later slide"
            }
        ],
        "source_date_issues": [],
        "numerical_inconsistencies": [],
        "cross_reference_issues": []
    }
    
    print(f"Running correction on {input_path}...")
    result = corrector.correct(
        original_path=input_path,
        validation_result=validation_result,
        output_path=output_path
    )
    
    print(f"Success: {result.success}")
    if not result.success:
        print(f"Error: {result.error_message}")
        
    print("Fixes applied:", result.fixes_applied)
    print("Fixes failed:", result.fixes_failed)
    
    # Check changes by slide
    print("Changes by slide:", result.changes_by_slide)

    # Verify output file content
    if result.success and Path(output_path).exists():
        print("\n--- Verifying Output PPTX ---")
        from pptx import Presentation
        prs = Presentation(output_path)
        for i, slide in enumerate(prs.slides, 1):
            print(f"Slide {i} shapes: {len(slide.shapes)}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "DEBUG" in text or "Benefit" in text:
                        print(f"  - Found Issue Box: '{text[:50]}...' at ({shape.left}, {shape.top})")

if __name__ == "__main__":
    debug_highlighting()
