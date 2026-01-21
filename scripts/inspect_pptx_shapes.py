
import sys
import os
from pptx import Presentation

def check_file(path):
    if not os.path.exists(path):
        print(f"File not found: {path}")
        return
        
    try:
        prs = Presentation(path)
        print(f"Checking {path}")
        print(f"Total slides: {len(prs.slides)}")
        
        total_shapes = 0
        box_count = 0
        for i, slide in enumerate(prs.slides, 1):
            shapes_on_slide = len(slide.shapes)
            total_shapes += shapes_on_slide
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and ("Suggested:" in shape.text or "Issue:" in shape.text):
                    box_count += 1
                    print(f"Slide {i}: Found Suggestion Box with text: {shape.text[:50]}...")

        print(f"Summary: Found {box_count} suggestion boxes across {total_shapes} shapes.")
    except Exception as e:
        print(f"Error reading file: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python scripts/inspect_pptx_shapes.py <path_to_pptx>")
    else:
        check_file(sys.argv[1])
