"""
Document Text Extractor
========================

This module focuses exclusively on **data extraction/structuring**, which is the
foundation for all subsequent compliance agents (disclaimer checks, ESG limits,
registration validation, etc.).  The goal is to produce a rich, opinionated
representation of a marketing document without performing any compliance logic
here – that work remains for later modules.

### Design Choices

* We rely on lightweight, locally executable libraries (`python-pptx`,
  `python-docx`, `pdf2image`, `pytesseract`) so the extraction step remains
  reproducible and auditable.  Cloud OCR services could offer higher accuracy,
  but they introduce latency, cost, and possibly data residency concerns.  We
  keep the interface flexible so a future implementation could swap the OCR
  layer behind the same data structure if needed.
* The extractor enriches the raw text with **presentation-aware metadata** such
  as slide roles, detected disclaimers, country mentions, and table source
  notes.  These heuristics are intentionally simple (regex/keyword driven) so
  they are deterministic and easy to troubleshoot.  More advanced NLP models
  can sit on top later, but the downstream compliance agents already have
  immediately useful signals.
* Where heuristics may produce false positives (e.g., any sentence containing a
  percentage might not be a performance statement), we surface both the signal
  and the original text fragment so that validators – human or automated – can
  double-check the context.

Every block below includes comments describing **why** we capture certain
fields, and why we prefer these pragmatic heuristics over more complex
approaches at this stage.
"""

import os
import re
from pathlib import Path
from typing import Dict, List, Any, Tuple, Optional
from io import BytesIO

# fama librairies li nest3mlouhom fil extraction manest7a9ouch OCR w image processing 
# Import extraction libraries (required). Missing packages will raise ImportError.
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docx import Document as DocxDocument

import pytesseract

from pdf2image import convert_from_path

from PIL import Image, ImageOps, ImageFilter

from langdetect import detect as detect_language


# ---------------------------------------------------------------------------
# Domain-driven keyword lists (kept strategically simple so they remain easy to
# audit against the Consignes document).  Whenever downstream teams refine the
# taxonomy we can centralize it here.
# ---------------------------------------------------------------------------

# Legal disclaimers often include these anchors.  We prefer explicit keywords
# rather than fuzzy similarity so that false positives are easy to spot and
# update.
DISCLAIMER_KEYWORDS = {
    "disclaimer",
    "legal notice",
    "risk warning",
    "important information",
    "performance disclaimer",
    "back-tested performance",
    "simulation of future performance",
}

# Glossary sections have a distinctive label in the template.  Again, using a
# simple keyword avoids accidentally classifying unrelated slides.
GLOSSARY_KEYWORDS = {"glossary", "definitions"}

# Countries referenced by the Registration agent.  Mirrors the list used in the
# feature extractor so rules rely on the same vocabulary.
COUNTRY_NAMES = {
    "germany", "austria", "belgium", "chile", "spain", "france",
    "italy", "luxembourg", "netherlands", "peru", "portugal",
    "united kingdom", "uk", "singapore", "sweden", "switzerland",
    "finland", "denmark", "norway", "ireland", "united arab emirates",
    "uae", "iceland"
}

COUNTRY_SECTION_KEYWORDS = {
    "distribution": "distribution",
    "available": "available",
    "enregistr": "enregistr",  # captures enregistré/enregistrée
    "registr": "registr",
    "marketing": "commercialisation",
}

ISSUER_SECTION_KEYWORDS = {
    "top holdings",
    "principales positions",
    "top 10",
    "top issuers",
    "main holdings",
    "largest positions",
    "top positions",
}

LEGAL_NOTICE_KEYWORDS = [
    "oddo bhf asset management",
    "société de gestion",
    "management company",
    "regulated by",
    "autorité des marchés financiers",
]

# Percentage detection powers performance heuristics.  The regex accepts both
# comma and dot decimals because source files mix locales (e.g., "15,3%" vs
# "15.3%").
PERCENT_PATTERN = re.compile(r"\b\d{1,3}(?:[\.,]\d+)?\s?%")

# Date extraction for sources: we look for patterns like "Jan 2024", "01/2024",
# "2024-01-31".  This remains intentionally permissive; downstream agents can
# decide if the granularity is sufficient.
DATE_PATTERN = re.compile(
    r"(\b\d{1,2}[\-/]\d{1,2}[\-/]\d{2,4}\b|\b\d{4}[\-/]\d{1,2}[\-/]\d{1,2}\b|\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}\b)",
    re.IGNORECASE,
)

# Simple PII detectors to warn if marketing material accidentally carries email
# addresses or phone numbers; compliance teams flagged this as a risk.
EMAIL_PATTERN = re.compile(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}")
PHONE_PATTERN = re.compile(r"\+?\d[\d\s().-]{6,}")

ISIN_PATTERN = re.compile(r"\b[A-Z]{2}[A-Z0-9]{9}[0-9]\b")
SHARE_CLASS_PATTERN = re.compile(r"\b[A-Z]{1,4}-[A-Z]{3}\b")

DISCLAIMER_CATEGORIES = {
    # Mandatory disclaimers
    'obam_presentation': ["oddo bhf asset management", "obam", "oddo bhf am", "oddo bhf asset management sas", "oddo bhf am gmbh", "oddo bhf am lux"],
    'commercial_documentation': ["commercial documentation", "documentation commerciale", "commerciale dokumentation"],
    'commercial_documentation_raif': ["raif", "luxembourg funds-raif", "luxembourg funds raif", "fonds luxembourg-raif"],
    
    # Content-specific disclaimers
    'performance': ["past performance", "performances passées", "performance disclaimer", "past performance is not", "performances passées ne préjugent pas"],
    'esg_risk': ["esg", "sustainable", "article 8", "article 9", "article 6", "sustainability", "durabilité"],
    'issuers': ["issuer", "issuers", "emetteur", "émetteurs", "company logo", "logo", "entreprise mentionnée"],
    'simulation': ["simulation", "simulé", "forward-looking", "future performance", "projection", "prévision"],
    'backtest': ["back-tested", "backtested", "données arrière", "historique reconstruit", "back test", "back-test"],
    'new_offer': ["new product", "nouvelle offre", "nouvelle stratégie", "new strategy", "nouveau produit", "new offer"],
    'ytm': ["ytm", "ytw", "yield to maturity", "yield to worst", "yield-to-maturity", "yield-to-worst"],
    'opinion': ["opinion", "avis", "view", "vue", "forecast", "prévision"],
    'additional_information': ["additional information", "informations complémentaires", "zusätzliche informationen"],
    
    # Regulatory disclaimers
    'sfdr_article_6': ["article 6", "sfdr article 6", "article 6 sfdr"],
    'sfdr_article_8': ["article 8", "sfdr article 8", "article 8 sfdr"],
    'sfdr_article_9': ["article 9", "sfdr article 9", "article 9 sfdr"],
    'sfdr': ["sfdr", "sustainable finance disclosure regulation"],
    'money_market_fund': ["money market fund", "fonds monétaire", "geldmarktfonds", "mmf"],
    'sri': ["sri", "summary risk indicator", "indicateur de risque synthétique", "risikoindikator", "risk indicator"],
}


class DocumentExtractor:
    """
    Extract text and tables from various document formats
    """
    
    def __init__(self, enable_chart_analysis: bool = True):
        # Houni ninitializiw l'extractor wel fonction mtaa kol format 
        
        # Define supported formats and functions to call
        self.supported_formats = {
            '.pptx': self._extract_pptx,  # PowerPoint
            '.docx': self._extract_docx,  # Word
            '.pdf': self._extract_pdf     # PDF
        }
        self._configure_tesseract()
        
        # Initialize chart analyzer if enabled
        self.enable_chart_analysis = enable_chart_analysis
        self.chart_analyzer = None
        if enable_chart_analysis:
            try:
                from .chart_analyzer import ChartAnalyzer
                self.chart_analyzer = ChartAnalyzer(use_llm=True)
            except Exception as e:
                print(f"Warning: Chart analyzer not available: {e}")
                self.enable_chart_analysis = False

    def _configure_tesseract(self) -> None:
        """Configure Tesseract path from environment variable if provided."""
        if pytesseract is None:
            return

        tesseract_cmd = os.getenv("TESSERACT_CMD")
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        self.tesseract_lang = os.getenv("TESSERACT_LANG", "eng")
        self.tesseract_config = os.getenv("TESSERACT_CONFIG", "--psm 6")

    # ------------------------------------------------------------------
    # Utility helpers (kept inside the class so they can use the same keyword
    # dictionaries and be overridden in tests if needed).
    # ------------------------------------------------------------------

    def _split_sentences(self, text: str) -> List[str]:
        # We intentionally apply a conservative regex splitter instead of a
        # heavy NLP model.  The documents are often short bullet sentences where
        # rule-based splitting works well and keeps runtime low.
        return [s.strip() for s in re.split(r"(?<=[.!?])\s+", text) if s.strip()]

    def _detect_language(self, text: str) -> Optional[str]:
        # `langdetect` needs a few characters of context; if the snippet is too
        # short we skip detection to avoid misleading guesses.  We only run this
        # when the library is available – it’s optional but recommended.
        if detect_language is None or len(text) < 20:
            return None
        try:
            return detect_language(text)
        except Exception:
            return None

    def _detect_disclaimer_flags(self, text: str) -> bool:
        lowered = text.lower()
        return any(keyword in lowered for keyword in DISCLAIMER_KEYWORDS)

    def _detect_glossary(self, text: str) -> bool:
        lowered = text.lower()
        return any(keyword in lowered for keyword in GLOSSARY_KEYWORDS)

    def _detect_countries(self, text: str) -> List[str]:
        lowered = text.lower()
        matches = [country for country in COUNTRY_NAMES if country in lowered]
        # We return unique, title-cased country names to simplify downstream
        # comparison against registration datasets.
        return sorted({country.title() for country in matches})

    def _detect_performance_blocks(self, text: str) -> List[str]:
        blocks = []
        for sentence in self._split_sentences(text):
            if PERCENT_PATTERN.search(sentence):
                blocks.append(sentence)
        return blocks

    def _detect_source_notes(self, text: str) -> List[str]:
        lowered = text.lower()
        notes = []
        if "source" in lowered or "date" in lowered:
            # Source/date tagging is intentionally permissive: the compliance
            # agent would rather inspect a few extra sentences than miss the
            # only place where the source is mentioned.  We therefore capture
            # the entire text block and let downstream logic decide if it meets
            # the style requirements (e.g., bold, proximity to chart).
            notes.append(text.strip())
        return notes

    def _parse_table_source(self, note: str) -> Dict[str, Optional[str]]:
        # Extract a probable source name by removing date fragments and trailing
        # punctuation.  This is heuristic-based but provides structured fields
        # for the data-consistency agent without requiring manual annotation.
        date_match = DATE_PATTERN.search(note)
        date_text = date_match.group(0) if date_match else None
        source_candidate = note
        if date_text:
            source_candidate = source_candidate.replace(date_text, "")
        source_candidate = re.sub(r"^[^A-Za-z0-9]+|[^A-Za-z0-9]+$", "", source_candidate).strip()
        return {
            'raw_note': note,
            'source_name': source_candidate or None,
            'source_date': date_text,
        }

    def _extract_title_fields(self, text: str) -> Dict[str, Optional[str]]:
        """Extract key facts from the opening section of a document.

        The goal is not to be perfect, but to provide downstream compliance
        agents with a consistent set of cues (fund name, date, client type, etc.)
        so they do not need to re-tokenize the title slide.  We intentionally use
        simple regex heuristics so the behavior is deterministic and easy to
        tune when templates evolve.
        """
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        joined = " ".join(lines)

        # Fund name heuristic: first non-empty line with capital letters that is
        # not purely a date.  In the ODDO decks this reliably contains the fund
        # designation.
        fund_name = None
        for line in lines:
            if len(line) < 4:
                continue
            if DATE_PATTERN.search(line):
                continue
            if any(char.isalpha() for char in line):
                fund_name = line
                break

        # Currency: capture tokens like "EUR", "USD", "CHF" optionally wrapped
        # in parentheses.
        currency_match = re.search(r"(?:currency|devise)\s*[:\-]?\s*([A-Z]{3})", joined, re.IGNORECASE)
        if not currency_match:
            # Fallback: take an isolated 3-letter token that looks like an ISO
            # currency code but skip common ODDO acronyms.
            fallback_match = re.search(r"\b([A-Z]{3})\b", joined)
            if fallback_match and fallback_match.group(1) not in {"ODD", "ODDO", "BHF"}:
                currency_match = fallback_match
        currency = currency_match.group(1) if currency_match else None

        # Document date: reuse DATE_PATTERN for flexibility (supports "Jan 2024"
        # and "31/12/2024").
        date_match = DATE_PATTERN.search(joined)
        document_date = date_match.group(0) if date_match else None

        # Client type: check for explicit mentions.  Consignes differentiates
        # professional vs retail.
        client_type = None
        lowered = joined.lower()
        if "professional" in lowered or "professionnel" in lowered:
            client_type = "professional"
        elif "retail" in lowered or "non-professional" in lowered or "grand public" in lowered:
            client_type = "retail"

        # Risk indicator: look for SRRI or scale references.
        risk_indicator = None
        risk_match = re.search(r"srri\s*(\d)" , lowered)
        if risk_match:
            risk_indicator = risk_match.group(1)
        elif "risk" in lowered and any(token in lowered for token in ["1/7", "2/7", "3/7", "4/7", "5/7", "6/7", "7/7"]):
            risk_indicator = re.search(r"([1-7]/7)", lowered).group(1)  # type: ignore

        # Management company: explicit mention (ODDO BHF...).  Keep as raw text
        # so the compliance team can match exact wording from Glossaire.
        mgmt_match = re.search(r"ODDO\s+BHF[^\n]*", joined, re.IGNORECASE)
        management_company = mgmt_match.group(0) if mgmt_match else None

        return {
            'fund_name': fund_name,
            'currency': currency,
            'document_date': document_date,
            'client_type': client_type,
            'risk_indicator': risk_indicator,
            'management_company': management_company,
        }

    def _extract_identifiers(self, text: str, location: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Harvest ISINs and share-class codes with location metadata.

        We keep the detection logic straightforward so compliance reviewers can
        easily backtrack why a particular identifier was captured.
        """
        identifiers: List[Dict[str, Any]] = []
        for isin in ISIN_PATTERN.findall(text):
            entry = {
                'type': 'isin',
                'value': isin,
                'context': text.strip()[:120],
            }
            entry.update(location)
            identifiers.append(entry)

        for share in SHARE_CLASS_PATTERN.findall(text):
            entry = {
                'type': 'share_class',
                'value': share,
                'context': text.strip()[:120],
            }
            entry.update(location)
            identifiers.append(entry)

        return identifiers

    def _categorize_disclaimer(self, text: str) -> List[str]:
        categories = []
        lowered = text.lower()
        for category, keywords in DISCLAIMER_CATEGORIES.items():
            if any(keyword in lowered for keyword in keywords):
                categories.append(category)
        return categories

    def _country_entries(self, countries: List[str], heading: Optional[str], sentence: str, location: Dict[str, Any]) -> List[Dict[str, Any]]:
        entries = []
        for country in countries:
            entry = {
                'country': country,
                'heading': heading,
                'sentence': sentence.strip()[:160],
            }
            entry.update(location)
            entries.append(entry)
        return entries

    def _extract_issuer_mentions(self, text: str, location: Dict[str, Any]) -> List[Dict[str, Any]]:
        mentions = []
        lowered = text.lower()
        if not any(keyword in lowered for keyword in ISSUER_SECTION_KEYWORDS):
            return mentions

        # Split lines – bullet lists are common.  We treat any line with more
        # than one word and minimal digits as a potential issuer.
        for line in text.splitlines():
            candidate = line.strip(" •-\t")
            if len(candidate) < 3:
                continue
            if any(char.isdigit() for char in candidate):
                continue
            if candidate.lower() in ISSUER_SECTION_KEYWORDS:
                continue
            entry = {
                'issuer_name': candidate,
                'context': text.strip()[:160],
            }
            entry.update(location)
            mentions.append(entry)
        return mentions

    def _detect_legal_notice(self, text: str) -> bool:
        lowered = text.lower()
        return any(keyword in lowered for keyword in LEGAL_NOTICE_KEYWORDS)

    def _analyze_performance_sentence(self, sentence: str) -> Dict[str, Any]:
        lowered = sentence.lower()

        basis = None
        if any(keyword in lowered for keyword in ["net of fees", "net des frais", "net"]):
            basis = "net"
        elif any(keyword in lowered for keyword in ["gross of fees", "brut", "gross"]):
            basis = "gross"
        elif "backtest" in lowered or "back-tested" in lowered or "reconstitu" in lowered:
            basis = "backtest"
        elif "simulation" in lowered or "prospective" in lowered:
            basis = "simulation"

        benchmark_name = None
        benchmark_match = re.search(r"(?:benchmark|indice\s*(?:de\s*r\u00e9f\u00e9rence)?)[:\-]?\s*([^.;\n]+)", sentence, re.IGNORECASE)
        if benchmark_match:
            benchmark_name = benchmark_match.group(1).strip()

        period = None
        period_patterns = ["ytd", "1y", "3y", "5y", "10y", "since inception", "depuis l'origine"]
        for token in period_patterns:
            if token in lowered:
                period = token
                break

        value_float = None
        percent_match = PERCENT_PATTERN.search(sentence)
        if percent_match:
            percent_text = percent_match.group(0).replace('%', '').replace(' ', '').replace(',', '.')
            try:
                value_float = float(percent_text)
            except ValueError:
                value_float = None

        return {
            'sentence': sentence,
            'basis': basis,
            'benchmark': benchmark_name,
            'period': period,
            'value': value_float,
        }

    def _parse_numeric_value(self, cell: str) -> Optional[float]:
        cleaned = cell.strip().replace('%', '').replace(' ', '').replace(',', '.')
        if not cleaned:
            return None
        try:
            return float(cleaned)
        except ValueError:
            return None

    def _normalize_table_rows(self, rows: List[List[str]]) -> List[Dict[str, Any]]:
        if not rows or len(rows) < 2:
            return []
        headers = [str(h).strip() for h in rows[0]]
        entries = []
        for row in rows[1:]:
            if not row:
                continue
            label = str(row[0]).strip()
            if not label:
                continue
            for idx, cell in enumerate(row[1:], start=1):
                cell_text = str(cell).strip()
                if not cell_text:
                    continue
                value = self._parse_numeric_value(cell_text)
                if value is None:
                    continue
                column = headers[idx] if idx < len(headers) else f'column_{idx}'
                entries.append({
                    'label': label,
                    'column': column,
                    'value': value,
                    'raw': cell_text,
                })
        return entries
    
    def extract(self, file_path: str) -> Dict[str, Any]:
        # hadha bech ne5tarou bih extractor li yelzem based 3la lextension mta3 l fichier
        
        """
        Extract text and tables from document
        
        Args:
            
        Returns:
            Dict with:
        """
        file_path = Path(file_path)
        

        # Check if file exists
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        

        # Determine file type from extension
        file_ext = file_path.suffix.lower()
        
        # lahna lazem tkoun lformat mta3 file mawjouda fil supported_formats sinon mayet9rach
        # Check if format is supported
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_ext}")
        

        # Get appropriate extraction function
        extractor_func = self.supported_formats[file_ext]
        
        try:

            # Call appropriate function (pptx, docx, or pdf)
            result = extractor_func(file_path)
            result['extraction_method'] = file_ext[1:]
            result['file_path'] = str(file_path)
            result['file_size'] = file_path.stat().st_size
            return result
        except Exception as e:

            # If error occurs, return error
            return {
                'error': str(e),
                'extraction_method': file_ext[1:],
                'file_path': str(file_path)
            }
    
    def _extract_pptx(self, file_path: Path) -> Dict[str, Any]:
        # Yextracti text w tables men PowerPoint (slides, shapes, tables)
        """Extract text and tables from PowerPoint file"""
        if Presentation is None:
            raise ImportError("python-pptx not installed")
        
        prs = Presentation(file_path)
        
        full_text = []
        slides_data = []
        tables_data = []
        slide_summaries = []
        performance_sections = []
        disclaimer_slides = []
        glossary_slides = []
        country_mentions = set()
        table_sources = []
        identifiers: List[Dict[str, Any]] = []
        performance_table_entries: List[Dict[str, Any]] = []
        country_entries: List[Dict[str, Any]] = []
        issuer_mentions: List[Dict[str, Any]] = []
        legal_notice_slide = None
        charts_data: List[Dict[str, Any]] = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_ocr = []
            raw_source_notes = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_text.append(text)
                    full_text.append(text)
                    raw_source_notes.extend(self._detect_source_notes(text))
                
                # Extract tables - check shape type first
                table_data = None
                try:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    if (hasattr(shape, "shape_type") and 
                        shape.shape_type == MSO_SHAPE_TYPE.TABLE and 
                        hasattr(shape, "table")):
                        table_data = self._extract_pptx_table(shape.table)
                except (ImportError, AttributeError, Exception) as e:
                    # Skip if not a table or if extraction fails
                    pass
                
                if table_data:
                        table_index = len(tables_data) + 1
                        tables_data.append({
                            'slide_number': slide_idx,
                            'table_index': table_index,
                            'data': table_data
                        })
                        for entry in self._normalize_table_rows(table_data):
                            entry.update({'slide_number': slide_idx, 'table_index': table_index})
                            performance_table_entries.append(entry)
 
                # OCR and chart analysis on embedded images
                if (
                    MSO_SHAPE_TYPE is not None
                    and Image is not None
                    and hasattr(shape, "shape_type")
                    and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                ):
                    try:
                        image_bytes = shape.image.blob
                        
                        # Try chart analysis first (if enabled)
                        chart_data = None
                        if self.enable_chart_analysis and self.chart_analyzer:
                            try:
                                chart_data = self.chart_analyzer.extract_chart_data_from_image(
                                    image_bytes,
                                    location={'slide_number': slide_idx}
                                )
                                
                                # If it's a chart, extract data
                                if chart_data.get('is_chart'):
                                    # Add chart performance values to performance sections
                                    perf_values = chart_data.get('performance_values', [])
                                    if perf_values:
                                        if not any(s.get('slide_number') == slide_idx for s in performance_sections):
                                            performance_sections.append({
                                                'slide_number': slide_idx,
                                                'sentences': [],
                                                'entries': []
                                            })
                                        section = next(s for s in performance_sections if s.get('slide_number') == slide_idx)
                                        section['entries'].extend(perf_values)
                                    
                                    # Add chart source/date info
                                    source_info = chart_data.get('source_date_info', {})
                                    if source_info.get('has_source') or source_info.get('has_date'):
                                        raw_source_notes.append(
                                            f"Source: {source_info.get('source_text', '')} "
                                            f"Date: {source_info.get('date_text', '')}"
                                        )
                                    
                                    # Store chart data
                                    charts_data.append({
                                        'slide_number': slide_idx,
                                        'chart_type': chart_data.get('chart_type'),
                                        'chart_title': chart_data.get('chart_title'),
                                        'data_points': chart_data.get('data_points', []),
                                        'performance_values': chart_data.get('performance_values', []),
                                        'source_date_info': source_info,
                                        'confidence': chart_data.get('confidence', 0.0)
                                    })
                            except Exception as e:
                                # Chart analysis failed, continue with OCR
                                pass
                        
                        # Also do OCR for text extraction
                        if pytesseract is not None:
                            with Image.open(BytesIO(image_bytes)) as img:
                                img = self._prepare_image_for_ocr(img)
                                ocr_text = pytesseract.image_to_string(
                                    img,
                                    lang=getattr(self, "tesseract_lang", "eng"),
                                    config=getattr(self, "tesseract_config", "")
                                ).strip()
                                if ocr_text:
                                    slide_ocr.append(ocr_text)
                                    full_text.append(ocr_text)
                                    raw_source_notes.extend(self._detect_source_notes(ocr_text))
                    except Exception:
                        continue
            
            combined_text = "\n".join([*slide_text, *slide_ocr])
            slide_countries = self._detect_countries(combined_text)
            country_mentions.update(slide_countries)

            heading = slide_text[0] if slide_text else None
            country_entries.extend(
                self._country_entries(
                    slide_countries,
                    heading,
                    combined_text,
                    {'slide_number': slide_idx}
                )
            )

            issuer_mentions.extend(
                self._extract_issuer_mentions(
                    combined_text,
                    {'slide_number': slide_idx, 'heading': heading}
                )
            )

            identifier_entries = self._extract_identifiers(
                combined_text,
                {'slide_number': slide_idx}
            )
            if identifier_entries:
                identifiers.extend(identifier_entries)

            performance_sentences = self._detect_performance_blocks(combined_text)
            if performance_sentences:
                entries = [self._analyze_performance_sentence(sentence) for sentence in performance_sentences]
                performance_sections.append({
                    'slide_number': slide_idx,
                    'sentences': performance_sentences,
                    'entries': entries
                })

            if self._detect_disclaimer_flags(combined_text):
                disclaimer_slides.append(slide_idx)

            if self._detect_glossary(combined_text):
                glossary_slides.append(slide_idx)

            if legal_notice_slide is None and self._detect_legal_notice(combined_text):
                legal_notice_slide = slide_idx

            table_sources.extend({
                'slide_number': slide_idx,
                **self._parse_table_source(note)
            } for note in raw_source_notes)

            slides_data.append({
                'slide_number': slide_idx,
                'text': '\n'.join(slide_text),
                'ocr_text': '\n'.join(slide_ocr) if slide_ocr else None
            })

            slide_summaries.append({
                'slide_number': slide_idx,
                'is_title_candidate': slide_idx == 1,
                # We store booleans instead of high-level labels ("disclaimer",
                # "performance") so later agents can enforce more nuanced rules
                # without losing the raw detection signal.
                'contains_disclaimer': slide_idx in disclaimer_slides,
                'disclaimer_categories': self._categorize_disclaimer(combined_text) if slide_idx in disclaimer_slides else [],
                'contains_performance': bool(performance_sentences),
                'contains_glossary': slide_idx in glossary_slides,
                'countries': slide_countries,
                'language': self._detect_language(combined_text),
            })

            if slide_idx == 1:
                title_info = self._extract_title_fields(combined_text)

        structure = {
            'title_slide': 1 if slides_data else None,
            'disclaimer_slides': disclaimer_slides,
            'performance_slides': [section['slide_number'] for section in performance_sections],
            'glossary_slides': glossary_slides,
            'countries_detected': sorted(country_mentions),
            'disclaimer_categories': {
                slide['slide_number']: slide['disclaimer_categories']
                for slide in slide_summaries if slide['disclaimer_categories']
            },
            'country_entries': country_entries,
            'issuer_mentions': issuer_mentions,
            'has_glossary': bool(glossary_slides),
            'has_management_notice': legal_notice_slide is not None,
            'legal_notice_slide': legal_notice_slide,
        }

        return {
            'text': '\n'.join(full_text),
            'slides': slides_data,
            'slide_summaries': slide_summaries,
            'tables': tables_data,
            'total_slides': len(slides_data),
            'total_tables': len(tables_data),
            'table_sources': table_sources,
            'performance_sections': performance_sections,
            'performance_table_entries': performance_table_entries,
            'charts': charts_data,
            'total_charts': len(charts_data),
            'structure': structure,
            'identifiers': identifiers,
            'country_entries': country_entries,
            'issuer_mentions': issuer_mentions,
            'title_information': title_info if 'title_info' in locals() else None,
            'ocr': pytesseract is not None and any(slide.get('ocr_text') for slide in slides_data)
        }
    
    def _extract_pptx_table(self, table) -> List[List[str]]:
        # Yparsi table ta3 PPTX w yraja3 lista ta3 rows/cols
        """Extract data from PPTX table"""
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                row_data.append(cell_text)
            table_data.append(row_data)
        return table_data
    
    def _extract_docx(self, file_path: Path) -> Dict[str, Any]:
        # Yextracti paragraphs w tables men document Word (.docx)
        """Extract text and tables from Word document"""
        if DocxDocument is None:
            raise ImportError("python-docx not installed")
        
        doc = DocxDocument(file_path)
        
        full_text = []
        paragraphs_data = []
        tables_data = []
        paragraph_summaries = []
        performance_sections = []
        disclaimer_paragraphs = []
        glossary_paragraphs = []
        country_mentions = set()
        table_sources = []
        identifiers: List[Dict[str, Any]] = []
        country_entries: List[Dict[str, Any]] = []
        issuer_mentions: List[Dict[str, Any]] = []
        legal_notice_paragraph = None
        performance_table_entries: List[Dict[str, Any]] = []

        # Extract paragraphs
        for para_idx, paragraph in enumerate(doc.paragraphs, 1):
            text = paragraph.text.strip()
            if text:
                full_text.append(text)
                paragraphs_data.append({
                    'paragraph_number': para_idx,
                    'text': text,
                    'style': paragraph.style.name if paragraph.style else None
                })

                paragraph_countries = self._detect_countries(text)
                country_mentions.update(paragraph_countries)

                heading = paragraph.style.name if paragraph.style else None
                country_entries.extend(
                    self._country_entries(
                        paragraph_countries,
                        heading,
                        text,
                        {'paragraph_number': para_idx}
                    )
                )

                issuer_mentions.extend(
                    self._extract_issuer_mentions(
                        text,
                        {'paragraph_number': para_idx, 'heading': heading}
                    )
                )

                perf_sentences = self._detect_performance_blocks(text)
                if perf_sentences:
                    entries = [self._analyze_performance_sentence(sentence) for sentence in perf_sentences]
                    performance_sections.append({
                        'paragraph_number': para_idx,
                        'sentences': perf_sentences,
                        'entries': entries
                    })

                if self._detect_disclaimer_flags(text):
                    disclaimer_paragraphs.append(para_idx)

                if self._detect_glossary(text):
                    glossary_paragraphs.append(para_idx)

                table_sources.extend({
                    'paragraph_number': para_idx,
                    **self._parse_table_source(note)
                } for note in self._detect_source_notes(text))

                if para_idx == 1:
                    title_info = self._extract_title_fields(text)

                identifier_entries = self._extract_identifiers(
                    text,
                    {'paragraph_number': para_idx}
                )
                if identifier_entries:
                    identifiers.extend(identifier_entries)

                if legal_notice_paragraph is None and self._detect_legal_notice(text):
                    legal_notice_paragraph = para_idx

                paragraph_summaries.append({
                    'paragraph_number': para_idx,
                    'style': paragraph.style.name if paragraph.style else None,
                    # We intentionally keep these as simple flags so that later
                    # compliance agents can layer their own rule logic (e.g.,
                    # "slide 2 must contain a disclaimer" etc.).
                    'contains_disclaimer': para_idx in disclaimer_paragraphs,
                    'disclaimer_categories': self._categorize_disclaimer(text) if para_idx in disclaimer_paragraphs else [],
                    'contains_performance': bool(perf_sentences),
                    'contains_glossary': para_idx in glossary_paragraphs,
                    'countries': paragraph_countries,
                    'language': self._detect_language(text),
                })

        # Extract tables
        for table_idx, table in enumerate(doc.tables, 1):
            table_data = self._extract_docx_table(table)
            if table_data:
                tables_data.append({
                    'table_index': table_idx,
                    'data': table_data
                })
                for entry in self._normalize_table_rows(table_data):
                    entry.update({'table_index': table_idx})
                    performance_table_entries.append(entry)
         
        image_texts = self._extract_docx_images(doc)
        if image_texts:
            full_text.extend(image_texts)
         
        return {
            'text': '\n'.join(full_text),
            'paragraphs': paragraphs_data,
            'tables': tables_data,
            'image_ocr_text': image_texts if image_texts else None,
            'paragraph_summaries': paragraph_summaries,
            'performance_sections': performance_sections,
            'disclaimer_paragraphs': disclaimer_paragraphs,
            'glossary_paragraphs': glossary_paragraphs,
            'table_sources': table_sources,
            'structure': {
                'disclaimer_paragraphs': disclaimer_paragraphs,
                'performance_paragraphs': [section['paragraph_number'] for section in performance_sections],
                'glossary_paragraphs': glossary_paragraphs,
                'countries_detected': sorted(country_mentions),
                'disclaimer_categories': {
                    para['paragraph_number']: para['disclaimer_categories']
                    for para in paragraph_summaries if para['disclaimer_categories']
                },
                'country_entries': country_entries,
                'issuer_mentions': issuer_mentions,
                'has_glossary': bool(glossary_paragraphs),
                'has_management_notice': legal_notice_paragraph is not None,
                'legal_notice_paragraph': legal_notice_paragraph,
            },
            'identifiers': identifiers,
            'country_entries': country_entries,
            'issuer_mentions': issuer_mentions,
            'performance_table_entries': performance_table_entries,
            'title_information': title_info if 'title_info' in locals() else None,
            'total_paragraphs': len(paragraphs_data),
            'total_tables': len(tables_data),
            'ocr': bool(image_texts)
        }
 
    def _extract_docx_table(self, table) -> List[List[str]]:
        # Yparsi table men DOCX w yraja3 rows/cols as text
        """Extract data from DOCX table"""
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                row_data.append(cell_text)
            table_data.append(row_data)
        return table_data

    def _extract_docx_images(self, doc) -> List[str]:
        """Perform OCR on inline images within a DOCX document."""
        if pytesseract is None or Image is None:
            return []

        ocr_texts = []
        for rel in doc.part.related_parts.values():
            content_type = getattr(rel, "content_type", "")
            if not content_type.startswith("image/"):
                continue
            try:
                with Image.open(BytesIO(rel.blob)) as img:
                    img = self._prepare_image_for_ocr(img)
                    text = pytesseract.image_to_string(
                        img,
                        lang=getattr(self, "tesseract_lang", "eng"),
                        config=getattr(self, "tesseract_config", "")
                    ).strip()
                    if text:
                        ocr_texts.append(text)
            except Exception:
                continue
        return ocr_texts
    
    def _extract_pdf(self, file_path: Path) -> Dict[str, Any]:
        # OCR-based extraction using pdf2image + pytesseract
        """Extract text from PDF using OCR."""
        # Prefer native text extraction (no OCR) when possible. We try PyPDF2
        # first; if it's not available or yields little/no text, we fall back
        # to image-based OCR using pdf2image + pytesseract.
        use_ocr = True
        pages_data: List[Dict[str, Any]] = []
        full_text: List[str] = []
        page_summaries: List[Dict[str, Any]] = []
        performance_sections: List[Dict[str, Any]] = []
        disclaimer_pages: List[int] = []
        glossary_pages: List[int] = []
        country_mentions = set()
        identifiers: List[Dict[str, Any]] = []
        country_entries: List[Dict[str, Any]] = []
        issuer_mentions: List[Dict[str, Any]] = []
        legal_notice_page = None
        performance_table_entries: List[Dict[str, Any]] = []

        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(str(file_path))
            extracted_texts: List[str] = []
            for page in reader.pages:
                try:
                    t = page.extract_text() or ""
                except Exception:
                    t = ""
                extracted_texts.append(t.strip())

            total_chars = sum(len(t) for t in extracted_texts)
            if total_chars > 200:
                # Good quality text-based PDF; avoid OCR.
                use_ocr = False
                for page_num, text in enumerate(extracted_texts, 1):
                    pages_data.append({
                        'page_number': page_num,
                        'text': text,
                        'ocr_engine': None
                    })
                    if text:
                        full_text.append(text)
                        page_countries = self._detect_countries(text)
                        country_mentions.update(page_countries)

                        perf_sentences = self._detect_performance_blocks(text)
                        if perf_sentences:
                            entries = [self._analyze_performance_sentence(sentence) for sentence in perf_sentences]
                            performance_sections.append({
                                'page_number': page_num,
                                'sentences': perf_sentences,
                                'entries': entries
                            })

                        if self._detect_disclaimer_flags(text):
                            disclaimer_pages.append(page_num)

                        if self._detect_glossary(text):
                            glossary_pages.append(page_num)

                        if page_num == 1 and 'title_info' not in locals():
                            title_info = self._extract_title_fields(text)

                        identifier_entries = self._extract_identifiers(text, {'page_number': page_num})
                        if identifier_entries:
                            identifiers.extend(identifier_entries)

                        heading = None
                        country_entries.extend(
                            self._country_entries(page_countries, heading, text, {'page_number': page_num})
                        )

                        issuer_mentions.extend(
                            self._extract_issuer_mentions(text, {'page_number': page_num})
                        )

                        if legal_notice_page is None and self._detect_legal_notice(text):
                            legal_notice_page = page_num
            else:
                use_ocr = True
        except Exception:
            use_ocr = True

        if use_ocr:
            # Ensure pdf2image and pytesseract are available for OCR fallback
            if convert_from_path is None:
                raise ImportError("pdf2image not installed")
            if pytesseract is None:
                raise ImportError("pytesseract not installed")

            images = convert_from_path(str(file_path), dpi=300)
            for page_num, image in enumerate(images, 1):
                image = self._prepare_image_for_ocr(image)

                ocr_text = pytesseract.image_to_string(
                    image,
                    lang=getattr(self, "tesseract_lang", "eng"),
                    config=getattr(self, "tesseract_config", "")
                )
                cleaned_text = ocr_text.strip()
                page_countries: List[str] = []
                pages_data.append({
                    'page_number': page_num,
                    'text': cleaned_text,
                    'ocr_engine': 'pytesseract'
                })
                if cleaned_text:
                    full_text.append(cleaned_text)

                    page_countries = self._detect_countries(cleaned_text)
                    country_mentions.update(page_countries)

                    perf_sentences = self._detect_performance_blocks(cleaned_text)
                    if perf_sentences:
                        entries = [self._analyze_performance_sentence(sentence) for sentence in perf_sentences]
                        performance_sections.append({
                            'page_number': page_num,
                            'sentences': perf_sentences,
                            'entries': entries
                        })

                    if self._detect_disclaimer_flags(cleaned_text):
                        disclaimer_pages.append(page_num)

                    if self._detect_glossary(cleaned_text):
                        glossary_pages.append(page_num)

                    if page_num == 1 and 'title_info' not in locals():
                        title_info = self._extract_title_fields(cleaned_text)

                    identifier_entries = self._extract_identifiers(cleaned_text, {'page_number': page_num})
                    if identifier_entries:
                        identifiers.extend(identifier_entries)

                    heading = None  # PDFs lose heading structure; keep placeholder
                    country_entries.extend(
                        self._country_entries(page_countries, heading, cleaned_text, {'page_number': page_num})
                    )

                    issuer_mentions.extend(
                        self._extract_issuer_mentions(cleaned_text, {'page_number': page_num})
                    )

                    if legal_notice_page is None and self._detect_legal_notice(cleaned_text):
                        legal_notice_page = page_num

                page_summaries.append({
                    'page_number': page_num,
                    'contains_disclaimer': page_num in disclaimer_pages,
                    'disclaimer_categories': self._categorize_disclaimer(cleaned_text) if cleaned_text and page_num in disclaimer_pages else [],
                    'contains_performance': any(section['page_number'] == page_num for section in performance_sections),
                    'contains_glossary': page_num in glossary_pages,
                    'countries': page_countries if cleaned_text else [],
                    'language': self._detect_language(cleaned_text) if cleaned_text else None,
                })

        # If pages_data was populated by text-extraction branch, ensure page_summaries
        # for that branch are present (they are populated inline above), otherwise
        # page_summaries were filled during OCR loop.

        return {
            'text': '\n'.join(full_text),
            'pages': pages_data,
            'page_summaries': page_summaries,
            'performance_sections': performance_sections,
            'disclaimer_pages': disclaimer_pages,
            'glossary_pages': glossary_pages,
            'structure': {
                'disclaimer_pages': disclaimer_pages,
                'performance_pages': [section['page_number'] for section in performance_sections],
                'glossary_pages': glossary_pages,
                'countries_detected': sorted(country_mentions),
                'disclaimer_categories': {
                    page['page_number']: page['disclaimer_categories']
                    for page in page_summaries if page.get('disclaimer_categories')
                },
                'country_entries': country_entries,
                'issuer_mentions': issuer_mentions,
                'has_glossary': bool(glossary_pages),
                'has_management_notice': legal_notice_page is not None,
                'legal_notice_page': legal_notice_page,
            },
            'identifiers': identifiers,
            'country_entries': country_entries,
            'issuer_mentions': issuer_mentions,
            'performance_table_entries': performance_table_entries,
            'title_information': title_info if 'title_info' in locals() else None,
            'tables': [],
            'total_pages': len(pages_data),
            'total_tables': 0,
            'ocr': use_ocr
        }

    def _prepare_image_for_ocr(self, image: Image.Image) -> Image.Image:
        if image.mode != "RGB":
            image = image.convert("RGB")
        if ImageOps:
            image = ImageOps.grayscale(image)
            image = ImageOps.autocontrast(image)
        if ImageFilter:
            image = image.filter(ImageFilter.MedianFilter(size=3))
        return image


def extract_document(file_path: str) -> Dict[str, Any]:
    # yaayet l DocumentExtractor() w yextracti document
    """Convenience function to extract document"""
    extractor = DocumentExtractor()
    return extractor.extract(file_path)


if __name__ == "__main__":
    # hadhi juste ntastiw beha l'extractor
    # Nbadlou test_file b path mta3 document li n7ebou aalih
    test_file = "dataset/example_1/47861-6PG-FR-ODDO BHF Algo Trend US-20250831 v3 pn.pptx"
    
    if os.path.exists(test_file):
        print(f"Extracting from: {test_file}")
        extractor = DocumentExtractor()
        result = extractor.extract(test_file)
        
        print(f"\nExtraction Results:")
        print(f"Method: {result.get('extraction_method')}")
        print(f"Text length: {len(result.get('text', ''))} characters")
        print(f"Slides: {result.get('total_slides', 0)}")
        print(f"Tables: {result.get('total_tables', 0)}")
        
        # Show first 500 characters
        text = result.get('text', '')
        if text:
            print(f"\nFirst 500 characters:")
            print(text[:500])
    else:
        print(f"Test file not found: {test_file}")

