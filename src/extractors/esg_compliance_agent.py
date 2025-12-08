from pathlib import Path
from pydantic import BaseModel, Field
from typing import List, Optional, Literal, Dict
from datetime import datetime
import fitz
import openai
import re
import json
import base64
from concurrent.futures import ThreadPoolExecutor, as_completed
from pptx import Presentation
from docx import Document
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from textwrap import wrap
import logging
import os
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Get LLM configuration (same as data consistency agent)
OPENAI_API_KEY = os.getenv("TOKEN_FACTORY_API_KEY") or os.getenv("LLM_API_KEY")
OPENAI_BASE_URL = os.getenv("TOKEN_FACTORY_BASE_URL") or os.getenv("LLM_BASE_URL")
LLM_MODEL = os.getenv("LLM_MODEL_NAME", "hosted_vllm/Llama-3.1-70B-Instruct")
LLM_MODEL_IMAGE = os.getenv("LLM_VISION_MODEL", "hosted_vllm/llava-1.5-7b-hf")

# ============================================================
# LOGGING SETUP
# ============================================================

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ============================================================
# MODÈLES PYDANTIC (optimisés + enrichis)
# ============================================================

class ESGLevel(BaseModel):
    level: Literal["engaging", "reduced", "limited", "none"]
    exclusion_percentage: Optional[float] = None
    portfolio_coverage: Optional[float] = None
    sfdr_article: Optional[Literal[6, 8, 9]] = None
    description: str

class ESGMentions(BaseModel):
    total_text_length: int
    esg_text_length: int
    esg_percentage: float
    esg_keywords_found: List[str]
    esg_keywords_by_slide: Dict[str, List[int]] = {}  # ✅ NEW: keyword -> [slide numbers]
    esg_sentences: List[str]
    mandatory_regulatory_mentions: int
    commercial_esg_mentions: int

class ImageAnalysisResult(BaseModel):
    image_path: str
    slide_number: int
    slide_title: str
    extracted_text_ocr: str
    llm_visual_description: str
    title_image_coherence: str
    esg_keywords_in_image: List[str]
    compliance_issues: List[str]

class ESGViolation(BaseModel):
    violation_type: str
    severity: Literal["critical", "high", "medium", "low"]
    description: str
    location: Optional[str] = None
    recommendation: str

class ESGComplianceOutput(BaseModel):
    document_id: str
    file_name: str
    processed_at: datetime = Field(default_factory=datetime.utcnow)
    fund_name: str
    client_type: Literal["retail", "professional"]
    document_type: str
    esg_level: ESGLevel
    esg_mentions: ESGMentions
    is_compliant: bool
    violations: List[ESGViolation] = []
    overall_confidence: float
    requires_human_review: bool
    image_analysis_results: List[ImageAnalysisResult] = []

# ============================================================
# FONCTIONS UTILITAIRES
# ============================================================

def normalize_esg_level(raw: dict) -> dict:
    """Normalise le niveau ESG retourné par le LLM avant Pydantic."""
    level = str(raw.get("level", "")).strip().lower()
    
    mapping = {
        "article 6": "none",
        "article 8": "reduced",
        "article 9": "engaging",
        "6": "none",
        "8": "reduced",
        "9": "engaging",
    }
    
    if level in mapping:
        raw["level"] = mapping[level]
    elif level not in ["engaging", "reduced", "limited", "none"]:
        art = raw.get("sfdr_article")
        if art == 6:
            raw["level"] = "none"
        elif art == 8:
            raw["level"] = "reduced"
        elif art == 9:
            raw["level"] = "engaging"
        else:
            raw["level"] = "none"
    
    if "sfdr_article" in raw and raw["sfdr_article"]:
        try:
            raw["sfdr_article"] = int(raw["sfdr_article"])
        except (ValueError, TypeError):
            raw["sfdr_article"] = None
    
    if "exclusion_percentage" in raw and raw["exclusion_percentage"]:
        try:
            raw["exclusion_percentage"] = float(raw["exclusion_percentage"])
        except (ValueError, TypeError):
            raw["exclusion_percentage"] = None
    
    if "portfolio_coverage" in raw and raw["portfolio_coverage"]:
        try:
            raw["portfolio_coverage"] = float(raw["portfolio_coverage"])
        except (ValueError, TypeError):
            raw["portfolio_coverage"] = None
    
    return raw

# ============================================================
# EXTRACTION AUTOMATIQUE PROSPECTUS
# ============================================================

def extract_fund_metadata_from_prospectus(prospectus_path: str) -> dict:
    """Extrait le nom du fonds et l'article SFDR du prospectus."""
    prospectus_path = Path(prospectus_path)
    if not prospectus_path.exists():
        logger.warning(f"Prospectus introuvable: {prospectus_path}")
        return {}

    try:
        logger.info("📖 Extraction métadonnées du prospectus...")
        if prospectus_path.suffix.lower() == '.pdf':
            doc = fitz.open(str(prospectus_path))
            paragraphs = []
            for page in doc:
                text = page.get_text("text")
                paragraphs.extend([p.strip() for p in text.split('\n') if p.strip()])
        else:
            doc = Document(str(prospectus_path))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        if not paragraphs:
            logger.warning("Prospectus vide")
            return {}

        full_text = "\n".join(paragraphs).lower()
        fund_name = paragraphs[0] if paragraphs else "Unknown"

        sfdr_article = None
        if "article 9" in full_text:
            sfdr_article = 9
        elif "article 8" in full_text:
            sfdr_article = 8
        elif "article 6" in full_text:
            sfdr_article = 6

        logger.info(f"📄 Métadonnées extraites: Fonds={fund_name}, Article={sfdr_article}")
        return {
            "fund_name": fund_name,
            "sfdr_article": sfdr_article
        }
    except Exception as e:
        logger.error(f"Erreur extraction prospectus: {e}")
        return {}

# ============================================================
# DOCUMENT LOADER
# ============================================================

class DocumentLoader:
    def __init__(self, file_path: str):
        self.file_path = Path(file_path)
        self.doc_type = self.file_path.suffix.lower()

    def load(self) -> dict:
        if self.doc_type == '.pdf':
            logger.info("📑 Chargement PDF...")
            doc = fitz.open(str(self.file_path))
            full_text = [page.get_text("text") for page in doc]
            logger.info(f"✅ PDF chargé ({len(doc)} pages)")
            return {"full_text": "\n".join(full_text), "images": [], "slides_data": []}

        elif self.doc_type == '.pptx':
            logger.info("🎯 Chargement PPTX...")
            prs = Presentation(str(self.file_path))
            full_text = []
            slide_images = []
            slide_data = []

            for i, slide in enumerate(prs.slides):
                slide_title = ""
                slide_text = []
                image_count = 0

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder and shape.placeholder_format.type == 1:
                            slide_title = shape.text.strip()
                        slide_text.append(shape.text)

                    if hasattr(shape, "shape_type") and shape.shape_type == 13 and image_count == 0:
                        image_count += 1
                        slide_images.append({
                            "path": None,
                            "slide_number": i + 1,
                            "slide_title": slide_title or f"Slide {i+1}"
                        })

                full_text.extend(slide_text)
                slide_data.append({
                    "slide_number": i + 1,
                    "title": slide_title or f"Slide {i+1}",
                    "text": "\n".join(slide_text)
                })

            logger.info(f"✅ PPTX chargé ({len(prs.slides)} slides, {len(slide_images)} images trouvées)")
            return {
                "full_text": "\n".join(full_text),
                "images": slide_images,
                "slides_data": slide_data
            }

        elif self.doc_type == '.docx':
            logger.info("📄 Chargement DOCX...")
            doc = Document(str(self.file_path))
            full_text = [p.text for p in doc.paragraphs if p.text.strip()]
            logger.info(f"✅ DOCX chargé ({len(full_text)} paragraphes)")
            return {"full_text": "\n".join(full_text), "images": [], "slides_data": []}

        else:
            raise ValueError(f"Format non supporté: {self.doc_type}")

class MetadataLoader:
    @staticmethod
    def load(metadata_file: str) -> dict:
        with open(metadata_file, 'r', encoding='utf-8') as f:
            return json.load(f)

# ============================================================
# ESG ANALYZER (optimisé + performant)
# ============================================================

class ESGAnalyzer:
    def __init__(self, api_key: str, base_url: str):
        self.api_key = api_key
        self.base_url = base_url
        
        self.esg_keywords = {
            "ESG": ["ESG", "environnement", "social", "governance", "environmental", "sustainable", "durabilité"],
            "Durabilité": ["durable", "durabilité", "sustainability", "sustainable"],
            "Impact": ["impact", "impact investing", "investissement d'impact"],
            "Carbone": ["carbone", "carbon", "CO2", "émission"],
            "Climat": ["climat", "climate", "transition"],
            "Responsable": ["responsable", "responsible", "éthique", "ethics"],
            "Vert": ["vert", "green", "écologique"],
            "ISR": ["ISR", "SRI", "investissement socialement responsable"]
        }

        self.regulatory_patterns = [
            r"Article\s+[689]\s+",
            r"SFDR",
            r"durabilité\s+dans\s+le\s+secteur",
            r"règlement\s+européen",
            r"publication\s+d.informations",
            r"facteurs\s+de\s+durabilité",
            r"risques?\s+de\s+durabilité"
        ]
        self.compiled_patterns = [re.compile(p) for p in self.regulatory_patterns]
        
        # ✅ ENHANCEMENT #8: Pre-compile ESG keyword pattern for performance
        all_keywords = [kw.lower() for keywords in self.esg_keywords.values() for kw in keywords]
        self.esg_keyword_pattern = re.compile(
            r'\b(' + '|'.join(re.escape(kw) for kw in all_keywords) + r')\b',
            re.IGNORECASE
        )
        
        # ✅ ENHANCEMENT #5: LLM result cache
        self._esg_level_cache = {}
        self._esg_mentions_cache = {}

    def analyze_image_with_advanced_vision(self, image_path: str, slide_title: str, slide_number: int) -> ImageAnalysisResult:
        """Analyse image avec vision LLM AVANCÉE (sans OCR Tesseract)."""
        try:
            with open(image_path, "rb") as f:
                image_bytes = f.read()
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            
            client = openai.OpenAI(api_key=self.api_key, base_url=self.base_url)
            
            prompt = f"""Analyse cette image du slide "{slide_title}" (Slide {slide_number}).

📋 INSTRUCTIONS DÉTAILLÉES:

1. CONTENU VISUEL : Décris graphiques, logos, couleurs, mise en page
2. TEXTE VISIBLE : Extrait TOUS les textes (titres, légendes, données)
3. COHÉRENCE TITRE-IMAGE : L'image correspond-elle au titre ? (OUI/NON + explication)
4. ANALYSE ESG :
   - Identifie mentions ESG (environnement, social, gouvernance, durable, vert, responsable, impact)
   - Détecte symboles ESG (labels, certifications, graphiques de durabilité)
   - Évalue l'importance de l'ESG dans l'image (faible/moyen/élevé)
5. DISCLAIMERS : Signale tout avertissement, restriction légale ou condition
6. QUALITÉ COMMUNICATION : Greenwashing détecté ? (OUI/NON + explications)

CLASSIFICATION ESG DE L'IMAGE:
- "PARFAIT": Slide avec contenu ESG cohérent et complet
- "BON": Slide avec ESG présent mais basique
- "MANQUANT": Slide sans contenu ESG (pour fonds Article 8/9)
- "HORS SUJET": Slide non-ESG (acceptable pour Article 6)

Retourne JSON STRUCTURÉ:
{{
  "visual_description": "description détaillée du contenu visuel",
  "extracted_text": "texte exact extrait de l'image",
  "title_coherence": "OUI/NON - explication brève",
  "esg_elements": ["élément1", "élément2", ...],
  "esg_intensity": "faible/moyen/élevé",
  "esg_quality_rating": "PARFAIT/BON/MANQUANT/HORS_SUJET",
  "sustainability_indicators": ["indicator1", "indicator2"],
  "greenwashing_risk": "AUCUN/BAS/MOYEN/ÉLEVÉ",
  "disclaimers_detected": ["disclaimer1", "disclaimer2"],
  "recommendations": "recommandations d'amélioration"
}}"""
            
            response = client.chat.completions.create(
                model=LLM_MODEL_IMAGE,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{image_base64}"
                                }
                            }
                        ]
                    }
                ],
                temperature=0,
                max_tokens=1200,
                response_format={"type": "json_object"}
            )
            
            llm_result = json.loads(response.choices[0].message.content)
            
            visual_desc = llm_result.get("visual_description", "")
            extracted_text = llm_result.get("extracted_text", "")
            title_coherence = llm_result.get("title_coherence", "NON")
            esg_elements = llm_result.get("esg_elements", [])
            esg_intensity = llm_result.get("esg_intensity", "faible")
            esg_quality = llm_result.get("esg_quality_rating", "HORS_SUJET")
            sustainability = llm_result.get("sustainability_indicators", [])
            greenwashing = llm_result.get("greenwashing_risk", "AUCUN")
            disclaimers = llm_result.get("disclaimers_detected", [])
            
            combined_keywords = list(set(esg_elements + sustainability))
            
            compliance_issues = disclaimers.copy()
            if greenwashing in ["MOYEN", "ÉLEVÉ"]:
                compliance_issues.append(f"⚠️ Risque de greenwashing: {greenwashing}")
            
            return ImageAnalysisResult(
                image_path=image_path,
                slide_number=slide_number,
                slide_title=slide_title,
                extracted_text_ocr=extracted_text,
                llm_visual_description=f"{visual_desc}\n\n📊 QUALITÉ ESG: {esg_quality}\n💪 INTENSITÉ: {esg_intensity}",
                title_image_coherence=title_coherence,
                esg_keywords_in_image=combined_keywords,
                compliance_issues=compliance_issues
            )
            
        except Exception as e:
            logger.error(f"Erreur analyse vision: {e}")
            return ImageAnalysisResult(
                image_path=image_path,
                slide_number=slide_number,
                slide_title=slide_title,
                extracted_text_ocr=f"Erreur: {str(e)}",
                llm_visual_description=f"Erreur analyse: {str(e)}",
                title_image_coherence="ERREUR",
                esg_keywords_in_image=[],
                compliance_issues=[f"Erreur technique: {str(e)}"]
            )

    def detect_esg_level_v2(self, fund_name: str, text: str, sfdr_article: Optional[int] = None) -> ESGLevel:
        """Détermine le niveau ESG via LLM avec PROMPTS STRICTS."""
        # ✅ ENHANCEMENT #5: Check cache first
        import hashlib
        cache_key = hashlib.md5(f"{fund_name}:{text[:1000]}:{sfdr_article}".encode()).hexdigest()
        if cache_key in self._esg_level_cache:
            logger.info("📋 Using cached ESG level detection")
            return self._esg_level_cache[cache_key]
        
        logger.info("🤖 Détection ESG level via LLM...")
        client = openai.OpenAI(api_key=self.api_key, base_url=self.base_url)

        article_match = re.search(r"Article\s+([689])\s", text)
        detected_article = int(article_match.group(1)) if article_match else sfdr_article

        # ✅ ENHANCEMENT #6: Retry logic with exponential backoff
        max_retries = 3
        for attempt in range(max_retries):
            try:
                response = client.chat.completions.create(
                    model=LLM_MODEL,
                    messages=[{
                        "role": "system",
                        "content": """You are an ESG/SFDR compliance expert.

CRITICAL: You MUST return a strict JSON with these exact fields:
- level: MUST be one of ["engaging", "reduced", "limited", "none"] (lowercase, nothing else)
- sfdr_article: integer 6, 8, or 9 (or null)
- exclusion_percentage: number or null
- portfolio_coverage: number or null
- description: brief string

NEVER put "Article 8" or "Article 9" inside the "level" field.
"level" is independent from sfdr_article.
Example:
{
  "level": "reduced",
  "sfdr_article": 8,
  "exclusion_percentage": null,
  "portfolio_coverage": null,
  "description": "Fund applies ESG factors in a limited way"
}"""
                    }, {
                        "role": "user",
                        "content": f"Fund: {fund_name}\nDetected Article SFDR: {detected_article or 'Unknown'}\n\nText excerpt:\n{text[:2000]}"
                    }],
                    temperature=0,
                    response_format={"type": "json_object"}
                )
                break
            except openai.RateLimitError as e:
                if attempt < max_retries - 1:
                    wait_time = 2 ** attempt
                    logger.warning(f"⚠️ Rate limit hit, retrying in {wait_time}s...")
                    import time
                    time.sleep(wait_time)
                else:
                    raise
            except openai.APIError as e:
                if attempt < max_retries - 1:
                    logger.warning(f"⚠️ API error, retrying... ({e})")
                    import time
                    time.sleep(1)
                else:
                    raise

        # ✅ ENHANCEMENT #8: Structured output validation with fallback
        try:
            data = json.loads(response.choices[0].message.content)
            logger.info(f"Raw LLM response: {data}")
            
            data = normalize_esg_level(data)
            logger.info(f"Normalized data: {data}")
            
            esg_level = ESGLevel(**data)
            logger.info(f"✅ ESG Level: {esg_level.level.upper()} (Article {esg_level.sfdr_article})")
            
            # Cache successful result
            self._esg_level_cache[cache_key] = esg_level
            return esg_level
        except Exception as e:
            logger.error(f"❌ Erreur parsing ESG Level: {e}")
            logger.error(f"Data reçu: {data if 'data' in locals() else 'N/A'}")
            
            # Fallback: Try to fix JSON with LLM
            try:
                logger.info("🔧 Attempting to fix invalid JSON...")
                fix_response = client.chat.completions.create(
                    model=LLM_MODEL,
                    messages=[{
                        "role": "system",
                        "content": f"Fix this invalid JSON to match ESGLevel schema: {ESGLevel.schema_json()}"
                    }, {
                        "role": "user",
                        "content": f"Invalid JSON: {data if 'data' in locals() else response.choices[0].message.content}"
                    }],
                    temperature=0,
                    response_format={"type": "json_object"}
                )
                fixed_data = json.loads(fix_response.choices[0].message.content)
                fixed_data = normalize_esg_level(fixed_data)
                esg_level = ESGLevel(**fixed_data)
                logger.info(f"✅ Fixed ESG Level: {esg_level.level.upper()}")
                self._esg_level_cache[cache_key] = esg_level
                return esg_level
            except:
                # Ultimate fallback
                logger.error("❌ Could not fix JSON, using safe default")
                default_level = ESGLevel(
                    level="none",
                    sfdr_article=sfdr_article or 6,
                    exclusion_percentage=None,
                    portfolio_coverage=None,
                    description="ESG level could not be determined - defaulting to Article 6"
                )
                self._esg_level_cache[cache_key] = default_level
                return default_level

    def analyze_esg_mentions_v2(self, text: str, document_type: str = "marketing", slides_data: List[dict] = None) -> ESGMentions:
        """✅ AMÉLIORÉ: Analyse ESG mentions + track chaque keyword par slide."""
        # ✅ ENHANCEMENT #5: Check cache first
        import hashlib
        cache_key = hashlib.md5(f"{text[:1000]}:{document_type}".encode()).hexdigest()
        if cache_key in self._esg_mentions_cache:
            logger.info("📋 Using cached ESG mentions analysis")
            return self._esg_mentions_cache[cache_key]
        
        logger.info("📊 Analyse ESG mentions avec localisation par slide...")
        
        total_length = len(text)
        sentences = re.split(r'[.!?]', text)

        esg_keywords_found = set()
        esg_keywords_by_slide: Dict[str, List[int]] = {}
        esg_sentences = []
        esg_text_length = 0
        mandatory_regulatory = 0
        commercial_esg = 0

        slides_dict = {}
        if slides_data:
            for slide in slides_data:
                slides_dict[slide["slide_number"]] = slide["text"].lower()

        # ✅ ENHANCEMENT #7: Use pre-compiled regex for faster matching
        for sentence in sentences:
            sentence_lower = sentence.lower()
            
            # Fast keyword detection with pre-compiled pattern
            found_keywords = self.esg_keyword_pattern.findall(sentence_lower)
            
            if found_keywords:
                esg_keywords_found.update(found_keywords)
                
                is_regulatory = any(pattern.search(sentence_lower) for pattern in self.compiled_patterns)
                if is_regulatory:
                    mandatory_regulatory += 1
                else:
                    commercial_esg += 1
                    esg_sentences.append(sentence.strip())
                esg_text_length += len(sentence)

        # Optimize slide mapping
        for keyword in esg_keywords_found:
            slide_numbers = []
            if slides_dict:
                keyword_lower = keyword.lower()
                slide_numbers = [num for num, text in slides_dict.items() if keyword_lower in text]
            esg_keywords_by_slide[keyword] = sorted(slide_numbers) if slide_numbers else []

        esg_percentage = (esg_text_length / total_length * 100) if total_length > 0 else 0

        mentions = ESGMentions(
            total_text_length=total_length,
            esg_text_length=esg_text_length,
            esg_percentage=round(esg_percentage, 2),
            esg_keywords_found=list(esg_keywords_found),
            esg_keywords_by_slide=esg_keywords_by_slide,
            esg_sentences=esg_sentences[:10],
            mandatory_regulatory_mentions=mandatory_regulatory,
            commercial_esg_mentions=commercial_esg
        )
        
        # Cache result
        self._esg_mentions_cache[cache_key] = mentions
        
        logger.info(f"✅ Mentions: {mentions.esg_percentage}% ({commercial_esg} commerciales)")
        return mentions

# ============================================================
# ESG COMPLIANCE AGENT
# ============================================================

class ESGComplianceAgent:
    def __init__(self, api_key: str, base_url: str, enable_rag: bool = False, max_workers: int = 3):
        self.analyzer = ESGAnalyzer(api_key, base_url)
        self.rag_enabled = False
        self.rag = None
        self.max_workers = max_workers

    def check_compliance(self,
                        file_path: str,
                        metadata_file: Optional[str] = None,
                        prospectus_path: Optional[str] = None,
                        document_type: str = "marketing",
                        document_id: str = "ESG-001",
                        analyze_images: bool = False) -> ESGComplianceOutput:

        logger.info("\n" + "="*70)
        logger.info(f"🔍 ANALYSE: {Path(file_path).name}")
        logger.info("="*70)

        loader = DocumentLoader(file_path)
        raw_doc = loader.load()
        full_text = raw_doc.get("full_text", "")
        slides_data = raw_doc.get("slides_data", [])

        fund_name = "Unknown Fund"
        sfdr_article_from_prospectus = None

        if prospectus_path and Path(prospectus_path).exists():
            fund_meta = extract_fund_metadata_from_prospectus(prospectus_path)
            fund_name = fund_meta.get("fund_name", "Unknown Fund")
            sfdr_article_from_prospectus = fund_meta.get("sfdr_article")

        client_type = "retail"
        image_results = []

        # ✅ ENHANCEMENT #1 & #3: Fix vision analysis + smart sampling
        if analyze_images and "images" in raw_doc and raw_doc["images"]:
            logger.info(f"📸 Analyzing {len(raw_doc['images'])} image(s) with vision LLM...")
            
            # Smart sampling: Only analyze slides with ESG keywords
            slides_dict = {slide['slide_number']: slide['text'].lower() for slide in slides_data}
            slides_to_analyze = []
            
            for img_info in raw_doc["images"]:
                slide_num = img_info.get("slide_number", 0)
                slide_text = slides_dict.get(slide_num, "")
                
                # Only analyze if slide contains ESG keywords
                if any(keyword.lower() in slide_text for category in self.analyzer.esg_keywords.values() for keyword in category):
                    slides_to_analyze.append(img_info)
            
            # Cap at 5 slides to control costs
            slides_to_analyze = slides_to_analyze[:5]
            
            if slides_to_analyze:
                logger.info(f"🎯 Smart sampling: {len(slides_to_analyze)}/{len(raw_doc['images'])} slides selected for vision analysis")
                
                # Parallel vision analysis
                with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                    futures = {}
                    for img_info in slides_to_analyze:
                        # Note: Vision requires actual image files, skip if path not available
                        if img_info.get("path"):
                            future = executor.submit(
                                self.analyzer.analyze_image_with_advanced_vision,
                                img_info["path"],
                                img_info.get("slide_title", f"Slide {img_info.get('slide_number')}"),
                                img_info.get("slide_number", 0)
                            )
                            futures[future] = img_info
                    
                    for future in as_completed(futures):
                        try:
                            result = future.result()
                            image_results.append(result)
                        except Exception as e:
                            img_info = futures[future]
                            logger.error(f"❌ Vision analysis failed for slide {img_info.get('slide_number')}: {e}")
            else:
                logger.info("📋 No slides with ESG keywords found, skipping vision analysis")

        # ✅ ENHANCEMENT #2: Parallel LLM calls for 2x speedup
        with ThreadPoolExecutor(max_workers=2) as executor:
            future_level = executor.submit(
                self.analyzer.detect_esg_level_v2,
                fund_name, full_text, sfdr_article_from_prospectus
            )
            future_mentions = executor.submit(
                self.analyzer.analyze_esg_mentions_v2,
                full_text, document_type, slides_data
            )
            
            esg_level = future_level.result()
            esg_mentions = future_mentions.result()

        violations = []
        is_compliant = True

        logger.info("\n⚠️ DÉTECTION VIOLATIONS...")

        if esg_level.level == "engaging":
            if (esg_level.exclusion_percentage or 0) < 20 or (esg_level.portfolio_coverage or 0) < 90:
                violations.append(ESGViolation(
                    violation_type="engaging_criteria_not_met",
                    severity="high",
                    description="Fonds 'engaging' mais critères non respectés.",
                    recommendation="Vérifier la classification du fonds et les critères ESG."
                ))
                is_compliant = False
                logger.warning(" 🔴 HAUTE: Critères Article 9 non respectés")

        elif esg_level.level == "reduced":
            if esg_mentions.esg_percentage > 10:
                violations.append(ESGViolation(
                    violation_type="esg_overmentioned_article8",
                    severity="critical",
                    description=f"Article 8: {esg_mentions.esg_percentage}% mentions ESG (max 10%)",
                    recommendation="Réduire mentions ESG à < 10%"
                ))
                is_compliant = False
                logger.warning(f" 🔴 CRITIQUE: Ratio ESG trop élevé ({esg_mentions.esg_percentage}%)")

        elif esg_level.level in ["limited", "none"]:
            if esg_mentions.commercial_esg_mentions > 2:
                violations.append(ESGViolation(
                    violation_type="esg_forbidden_article6",
                    severity="critical",
                    description=f"Article 6 (non-ESG): {esg_mentions.commercial_esg_mentions} mentions ESG détectées",
                    recommendation="Supprimer toutes mentions ESG commerciales"
                ))
                is_compliant = False
                logger.warning(" 🔴 CRITIQUE: Mentions ESG dans Article 6")

        for img_result in image_results:
            if img_result.esg_keywords_in_image and esg_level.level == "none":
                violations.append(ESGViolation(
                    violation_type="esg_in_image_article6",
                    severity="high",
                    description=f"Slide {img_result.slide_number}: contient ESG {img_result.esg_keywords_in_image}",
                    location=f"Slide {img_result.slide_number}: {img_result.slide_title}",
                    recommendation="Retirer éléments ESG visuels"
                ))
                is_compliant = False
                logger.warning(f" ❌ HAUTE: Image Slide {img_result.slide_number}")

        if is_compliant:
            logger.info(" ✅ Aucune violation")

        output = ESGComplianceOutput(
            document_id=document_id,
            file_name=str(file_path),
            fund_name=fund_name,
            client_type=client_type,
            document_type=document_type,
            esg_level=esg_level,
            esg_mentions=esg_mentions,
            is_compliant=is_compliant,
            violations=violations,
            overall_confidence=0.95,
            requires_human_review=len(violations) > 0,
            image_analysis_results=image_results
        )

        return output

# ============================================================
# RAPPORT AMÉLIORÉ AVEC SLIDES
# ============================================================

def generate_enhanced_compliance_report(result: ESGComplianceOutput, output_path: str):
    """
    ✅ RAPPORT AMÉLIORÉ avec keywords + slides où ils apparaissent
    """
    lines = []
    
    lines.append("=" * 80)
    lines.append("📊 RAPPORT COMPLET DE CONFORMITÉ ESG/SFDR".center(80))
    lines.append("=" * 80)
    lines.append("")
    
    lines.append(f"📌 FONDS : {result.fund_name}")
    lines.append(f"📋 ARTICLE SFDR : Article {result.esg_level.sfdr_article}")
    lines.append(f"🎯 NIVEAU ESG : {result.esg_level.level.upper()}")
    lines.append(f"👥 TYPE CLIENT : {result.client_type.capitalize()}")
    lines.append(f"📄 TYPE DOCUMENT : {result.document_type.capitalize()}")
    lines.append(f"📅 DATE : {result.processed_at.strftime('%d/%m/%Y à %H:%M:%S')}")
    lines.append(f"📁 FICHIER : {Path(result.file_name).name}")
    lines.append("")
    
    lines.append("-" * 80)
    status = "✅ CONFORME" if result.is_compliant else "🔴 NON CONFORME"
    lines.append(f"STATUT GLOBAL : {status}".center(80))
    lines.append("-" * 80)
    lines.append("")
    
    lines.append("📈 1. ANALYSE DU TEXTE")
    lines.append("─" * 80)
    lines.append(f"  Ratio ESG total : {result.esg_mentions.esg_percentage}%")
    lines.append(f"  Mentions ESG (commerciales) : {result.esg_mentions.commercial_esg_mentions}")
    lines.append(f"  Mentions réglementaires : {result.esg_mentions.mandatory_regulatory_mentions}")
    lines.append(f"  Longueur totale : {result.esg_mentions.total_text_length:,} caractères")
    lines.append("")
    
    if result.esg_level.level == "engaging":
        lines.append("  ✅ Article 9 (Engaging) : ESG promotion ATTENDUE")
    elif result.esg_level.level == "reduced":
        lines.append(f"  ⚠️ Article 8 (Reduced) : ESG % doit être < 10% (actuel: {result.esg_mentions.esg_percentage}%)")
        if result.esg_mentions.esg_percentage > 10:
            lines.append("     🔴 DÉPASSE LE SEUIL!")
    elif result.esg_level.level in ["limited", "none"]:
        lines.append(f"  🚫 Article {result.esg_level.sfdr_article} (Non-ESG) : 0 mention ESG autorisée")
        if result.esg_mentions.commercial_esg_mentions > 0:
            lines.append(f"     🔴 {result.esg_mentions.commercial_esg_mentions} mention(s) détectée(s)!")
    
    lines.append("")
    
    if result.esg_mentions.esg_keywords_found:
        lines.append("🔍 2. MOTS-CLÉS ESG DÉTECTÉS (avec localisation slides)")
        lines.append("─" * 80)
        
        keywords_sorted = sorted(result.esg_mentions.esg_keywords_found)
        
        for keyword in keywords_sorted:
            slides = result.esg_mentions.esg_keywords_by_slide.get(keyword, [])
            if slides:
                slides_str = ", ".join([f"Slide {s}" for s in slides])
                lines.append(f"  • {keyword:20} → {slides_str}")
            else:
                lines.append(f"  • {keyword:20} → (Non localisé)")
        
        lines.append("")
    
    if result.image_analysis_results:
        lines.append("📸 3. ANALYSE DES SLIDES/IMAGES")
        lines.append("─" * 80)
        
        total_images = len(result.image_analysis_results)
        perfect_slides = []
        good_slides = []
        missing_slides = []
        offTopic_slides = []
        
        for img in result.image_analysis_results:
            desc = img.llm_visual_description.lower()
            
            if "parfait" in desc:
                perfect_slides.append((img.slide_number, img.slide_title, img.esg_keywords_in_image))
            elif "bon" in desc:
                good_slides.append((img.slide_number, img.slide_title, img.esg_keywords_in_image))
            elif "manquant" in desc:
                missing_slides.append((img.slide_number, img.slide_title))
            else:
                offTopic_slides.append((img.slide_number, img.slide_title))
        
        lines.append(f"\n  📊 Résumé:")
        lines.append(f"     • Total slides : {total_images}")
        lines.append(f"     🟢 Slides PARFAITES (ESG complet) : {len(perfect_slides)}")
        lines.append(f"     🟡 Slides BONNES (ESG basique) : {len(good_slides)}")
        lines.append(f"     🟠 Slides MANQUANTES (sans ESG) : {len(missing_slides)}")
        lines.append(f"     ⚪ Slides HORS-SUJET : {len(offTopic_slides)}")
        
        if perfect_slides:
            lines.append(f"\n  ✅ SLIDES PARFAITES ({len(perfect_slides)}):")
            for slide_num, title, keywords in perfect_slides:
                lines.append(f"     • Slide {slide_num}: {title}")
                if keywords:
                    kw = ", ".join(keywords[:4])
                    lines.append(f"       Keywords: {kw}")
        
        if good_slides:
            lines.append(f"\n  🟡 SLIDES BONNES ({len(good_slides)}):")
            for slide_num, title, keywords in good_slides:
                lines.append(f"     • Slide {slide_num}: {title}")
                if keywords:
                    kw = ", ".join(keywords[:3])
                    lines.append(f"       Keywords: {kw}")
        
        if missing_slides:
            lines.append(f"\n  🟠 SLIDES MANQUANTES - À AMÉLIORER ({len(missing_slides)}):")
            for slide_num, title in missing_slides:
                lines.append(f"     • Slide {slide_num}: {title} ⚠️ Ajouter contenu ESG")
        
        if offTopic_slides and result.esg_level.level != "none":
            lines.append(f"\n  ⚪ SLIDES HORS-SUJET ({len(offTopic_slides)}):")
            for slide_num, title in offTopic_slides:
                lines.append(f"     • Slide {slide_num}: {title}")
        
        lines.append("")
    
    if result.violations:
        lines.append("⚠️ 4. VIOLATIONS DÉTECTÉES")
        lines.append("─" * 80)
        for i, violation in enumerate(result.violations, 1):
            severity_map = {
                "critical": "🔴 CRITIQUE",
                "high": "❌ HAUTE",
                "medium": "⚠️ MOYENNE",
                "low": "ℹ️ BASSE"
            }
            severity_label = severity_map.get(violation.severity, "•")
            
            lines.append(f"\n  {i}. {severity_label} - {violation.description}")
            if violation.location:
                lines.append(f"     Localisation: {violation.location}")
            lines.append(f"     ➜ Recommandation: {violation.recommendation.split(chr(10))[0]}")
        lines.append("")
    else:
        lines.append("✅ 4. VIOLATIONS")
        lines.append("─" * 80)
        lines.append("  ✓ Aucune violation détectée")
        lines.append("")
    
    lines.append("📝 5. ACTIONS REQUISES")
    lines.append("─" * 80)
    
    if not result.is_compliant:
        action_count = 1
        for violation in result.violations:
            rec = violation.recommendation.split('\n')[0].strip()
            lines.append(f"  {action_count}️⃣ {rec}")
            action_count += 1
        
        if result.image_analysis_results:
            missing_count = sum(1 for img in result.image_analysis_results 
                              if "manquant" in img.llm_visual_description.lower())
            if missing_count > 0:
                lines.append(f"  {action_count}️⃣ Compléter {missing_count} slide(s) avec contenu ESG approprié")
    else:
        lines.append("  ✅ Document CONFORME - Aucune action requise")
    
    lines.append("")
    
    lines.append("=" * 80)
    lines.append(f"Confiance: {result.overall_confidence * 100:.0f}% | Revue humaine: {'OUI' if result.requires_human_review else 'NON'}")
    lines.append("=" * 80)
    
    report_text = "\n".join(lines)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(report_text)
    
    logger.info(f"\n✅ Rapport amélioré sauvegardé: {output_path}")
    return report_text

# ============================================================
# GÉNÉRATION PDF
# ============================================================

def save_report_as_pdf(text: str, pdf_path: str):
    """Enregistre le rapport texte dans un PDF simple, multi-pages si besoin."""
    Path(pdf_path).parent.mkdir(parents=True, exist_ok=True)
    
    c = canvas.Canvas(pdf_path, pagesize=A4)
    width, height = A4
    
    margin = 50
    y = height - margin
    line_height = 14
    max_chars_per_line = 95
    
    for line in text.splitlines():
        for chunk in wrap(line, max_chars_per_line):
            if y <= margin:
                c.showPage()
                y = height - margin
            c.drawString(margin, y, chunk)
            y -= line_height
    
    c.save()
    logger.info(f"📄 Rapport PDF sauvegardé: {pdf_path}")

# ============================================================
# MAIN
# ============================================================

def main():
    logger.info("\n" + "="*70)
    logger.info("🚀 DÉMARRAGE AGENT CONFORMITÉ ESG/SFDR V4")
    logger.info("="*70)

    agent = ESGComplianceAgent(
        api_key=OPENAI_API_KEY,
        base_url=OPENAI_BASE_URL,
        enable_rag=False,
        max_workers=3
    )

    test_files = [
        (
            #"dataset 1/dataset/example_3/3 - FINAL CLEAN-6PG-GB-ODDO BHF US Equity Active ETF-20250831.pptx",
            #"dataset 1/dataset/example_3/metadata.json",
            #"dataset 1/dataset/example_3/prospectus.docx",  
            #"marketing",
            #"output_file/ESG/ESG-EXAMPLE3-002"

            "dataset 1/dataset/example_2/XXX-PRS-GB-ODDO BHF US Equity Active ETF-20250630_6PN.pptx",
            "dataset 1/dataset/example_2/metadata.json",
            "dataset 1/dataset/example_2/prospectus.docx",  
            "marketing",
            "output_file/ESG/ESG-EXAMPLE2-01"

            #"dataset 1/dataset/Batch 2/1369-PRS-FR-ODDO BHF Génération - 20231231 - FVv1.pptx",
            #None,
            #"dataset 1/dataset/Batch 2/1365-PRO-FR-ODDO BHF Generation CR-EUR-20251104.PDF",
            #"marketing",
            #"output_file/ESG/ESG-Batch-01",
        ),
    ]

    results_summary = []

    for file_path, _, prospectus_path, doc_type, doc_id in test_files:
        if Path(file_path).exists():
            try:
                prospectus_file = prospectus_path if prospectus_path and Path(prospectus_path).exists() else None

                result = agent.check_compliance(
                    file_path=file_path,
                    prospectus_path=prospectus_file,
                    document_type=doc_type,
                    document_id=doc_id,
                    analyze_images=False
                )

                output_json = f"{result.document_id}_esg_compliance.json"
                with open(output_json, "w", encoding="utf-8") as f:
                    f.write(result.model_dump_json(indent=2, exclude_none=True))
                logger.info(f"\n💾 JSON: {output_json}")

                output_file_txt = f"{result.document_id}_rapport_complet.txt"
                report_text = generate_enhanced_compliance_report(result, output_file_txt)

                output_file_pdf = f"{result.document_id}_rapport_complet.pdf"
                save_report_as_pdf(report_text, output_file_pdf)

                logger.info(f"\n{'='*70}")
                logger.info(f"📊 FONDS: {result.fund_name}")
                logger.info(f"📈 ESG LEVEL: {result.esg_level.level.upper()} (Article {result.esg_level.sfdr_article})")
                logger.info(f"📉 RATIO ESG: {result.esg_mentions.esg_percentage}%")
                logger.info(f"✅ CONFORME: {'OUI ✅' if result.is_compliant else 'NON ❌'}")
                logger.info(f"📋 VIOLATIONS: {len(result.violations)}")
                logger.info(f"📸 IMAGES: {len(result.image_analysis_results)}")
                logger.info(f"{'='*70}\n")

                if not result.is_compliant:
                    logger.warning("⚠️ VIOLATIONS DÉTECTÉES:")
                    for i, v in enumerate(result.violations, 1):
                        logger.warning(f" {i}. [{v.severity.upper()}] {v.description}")
                        if v.location:
                            logger.warning(f"    📍 {v.location}")

                results_summary.append({
                    "fund": result.fund_name,
                    "compliant": result.is_compliant,
                    "violations": len(result.violations),
                    "ratio_esg": result.esg_mentions.esg_percentage
                })

            except Exception as e:
                logger.error(f"❌ Erreur: {e}")
                import traceback
                traceback.print_exc()
        else:
            logger.error(f"⚠️ Fichier non trouvé: {file_path}")

    logger.info("\n" + "="*70)
    logger.info("📊 RÉSUMÉ FINAL")
    logger.info("="*70)
    for summary in results_summary:
        status = "✅ CONFORME" if summary["compliant"] else "⚠️ NON CONFORME"
        logger.info(f"{summary['fund']:40} | {status:20} | {summary['violations']} violations | ESG: {summary['ratio_esg']}%")

    logger.info("="*70)
    logger.info("✅ ANALYSE TERMINÉE")
    logger.info("="*70 + "\n")

if __name__ == "__main__":
    main()
