from pathlib import Path
from backend.extractors.core.document_extractor import DocumentExtractor

# Test with just first slide
pptx_path = Path("dataset/example_3/3 - FINAL CLEAN-6PG-GB-ODDO BHF US Equity Active ETF-20250831.pptx")

print(f"Testing chart analysis on first 3 images only...\n")

extractor = DocumentExtractor(enable_chart_analysis=True)

# Monkey patch to only process first slide
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(pptx_path)
image_count = 0

for slide_idx, slide in enumerate(prs.slides, 1):
    if slide_idx > 2:  # Only first 2 slides
        break
    print(f"\nSlide {slide_idx}:")
    for shape in slide.shapes:
        if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
            if image_count <= 3:  # Only first 3 images
                image_bytes = shape.image.blob
                result = extractor.chart_analyzer.extract_chart_data_from_image(
                    image_bytes,
                    location={'slide_number': slide_idx}
                )
                print(f"  Image {image_count}: is_chart={result.get('is_chart')}, confidence={result.get('confidence')}")

print(f"\n\nCheck c:\\temp\\chart_analysis_debug.log for full LLaVA responses")
