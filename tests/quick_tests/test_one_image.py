from pathlib import Path
from backend.extractors.core.chart_analyzer import ChartAnalyzer

# Clear cache first
from backend.utils.cache.llm_cache import get_llm_cache
cache = get_llm_cache()
if cache:
    print(f"Clearing LLM cache...")
    cache.clear()

# Now test ONE image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = Path("dataset/example_3/3 - FINAL CLEAN-6PG-GB-ODDO BHF US Equity Active ETF-20250831.pptx")
prs = Presentation(pptx_path)

analyzer = ChartAnalyzer(enable_caching=False)  # Disable caching for this test

print("\nTesting ONE image (slide 1, image 1) without cache...\n")

slide = list(prs.slides)[0]
for shape in slide.shapes:
    if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image_bytes = shape.image.blob
        print(f"Image size: {len(image_bytes)} bytes")
        
        result = analyzer.extract_chart_data_from_image(
            image_bytes,
            location={'slide_number': 1}
        )
        
        print(f"\nResult:")
        print(f"  is_chart: {result.get('is_chart')}")
        print(f"  confidence: {result.get('confidence')}")
        print(f"  chart_type: {result.get('chart_type')}")
        
        print(f"\n\nCheck c:\\temp\\chart_analysis_debug.log for LLaVA's full response")
        break  # Only test first image
