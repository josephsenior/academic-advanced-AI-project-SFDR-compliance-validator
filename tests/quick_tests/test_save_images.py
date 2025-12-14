from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path

pptx_path = Path("dataset/example_3/3 - FINAL CLEAN-6PG-GB-ODDO BHF US Equity Active ETF-20250831.pptx")
prs = Presentation(pptx_path)

output_dir = Path("c:/temp/saved_images")
output_dir.mkdir(exist_ok=True)

image_count = 0
for slide_idx, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
            image_bytes = shape.image.blob
            
            # Save first 3 images from each slide for inspection
            if image_count <= 10:
                output_path = output_dir / f"slide_{slide_idx}_image_{image_count}.png"
                with open(output_path, 'wb') as f:
                    f.write(image_bytes)
                print(f"Saved: {output_path.name} ({len(image_bytes)} bytes)")

print(f"\nTotal images: {image_count}")
print(f"Images saved to: {output_dir}")
